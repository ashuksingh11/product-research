import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_pptx(md_file, output_file):
    prs = Presentation()
    
    with open(md_file, 'r') as f:
        lines = f.readlines()
    
    title = "Market Intelligence Report: Kitchen AI 2026"
    subtitle = "Strategic Analysis & Competitive Landscape"
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    current_slide = None
    
    for line in lines:
        line = line.strip()
        if not line or line.startswith('---') or line.startswith('|'):
            continue
            
        if line.startswith('# '):
            slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(slide_layout)
            current_slide.shapes.title.text = line.replace('# ', '')
            tf = current_slide.shapes.placeholders[1].text_frame
            tf.word_wrap = True
            
        elif line.startswith('## '):
            if current_slide:
                p = tf.add_paragraph()
                p.text = line.replace('## ', '').upper()
                p.font.bold = True
                p.font.size = Pt(18)
            
        elif line.startswith('* ') or line.startswith('- '):
            if current_slide:
                p = tf.add_paragraph()
                p.text = line.strip('* -')
                p.level = 1
                p.font.size = Pt(14)
        else:
            if current_slide:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)

    prs.save(output_file)
    print(f"PPTX saved to {output_file}")

if __name__ == "__main__":
    create_pptx('comprehensive_market_report.md', 'comprehensive_market_report.pptx')
