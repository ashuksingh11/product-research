#!/usr/bin/env python3
"""
Samsung Kitchen: Comprehensive Product & Tech Competitive Landscape (Data-Heavy Variant)
4-slide VP/Director-level PPTX covering product features, AI/tech architecture, and strategy
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)
SAMSUNG_DARK = RGBColor(0x0A, 0x12, 0x5E)
SAMSUNG_LIGHT = RGBColor(0xE8, 0xEB, 0xF7)
LG_RED = RGBColor(0xA5, 0x00, 0x34)
BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
GE_BLUE = RGBColor(0x00, 0x4B, 0x8D)
WHIRLPOOL_GOLD = RGBColor(0xC8, 0x96, 0x23)
MIELE_BROWN = RGBColor(0x6E, 0x1D, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xF3, 0x9C, 0x12)
RED_ALERT = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=10, default_color=DARK_GRAY):
    """lines: list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


# ============================================================
# SLIDE 1: Samsung Product & Ecosystem Overview
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, WHITE)

# Top header bar
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Inches(0.9), SAMSUNG_BLUE)
add_text_box(slide1, Inches(0.5), Inches(0.15), Inches(9), Inches(0.35),
             "Samsung Kitchen: Product & AI Ecosystem Overview", 22, WHITE, True, font_name='Calibri')
add_text_box(slide1, Inches(0.5), Inches(0.52), Inches(9), Inches(0.3),
             "Comprehensive Product Research  |  March 2026  |  Confidential", 11, RGBColor(0xB0, 0xBE, 0xE8), False)

# Market context bar
add_shape(slide1, Inches(0), Inches(0.9), SLIDE_W, Inches(0.55), RGBColor(0xF0, 0xF2, 0xFA))
stats = [
    ("$311B", "Global Kitchen Market"),
    ("21.7%", "Samsung Smart Kitchen Share"),
    ("4.55/5", "Composite KPI Score (#1)"),
    ("425M+", "SmartThings Users"),
    ("9,304", "US Patents (2024)"),
]
stat_w = Inches(2.4)
for i, (val, label) in enumerate(stats):
    x = Inches(0.5) + i * stat_w
    add_text_box(slide1, x, Inches(0.92), stat_w, Inches(0.25), val, 14, SAMSUNG_BLUE, True)
    add_text_box(slide1, x, Inches(1.14), stat_w, Inches(0.25), label, 8, MED_GRAY, False)

# LEFT: Samsung Strengths (Product Features)
add_text_box(slide1, Inches(0.5), Inches(1.55), Inches(6), Inches(0.3),
             "SAMSUNG KEY PRODUCT DIFFERENTIATORS", 11, SAMSUNG_BLUE, True)

samsung_features = [
    ("Bespoke Design System", "8-14 color panels, MyBespoke custom options. Unmatched aesthetic\ncustomization at mass-premium scale. Strong in US and Korea.", SAMSUNG_BLUE),
    ("AI Family Hub+ (32\" Display)", "Largest touchscreen in industry. Google Gemini integration.\nAI Vision recognizes 1000s of items. Now Brief, FoodNote, streaming.", SAMSUNG_BLUE),
    ("AI Hybrid Cooling", "World-first compressor + Peltier dual cooling. Pork freshness 1.4x,\nsalmon 1.2x extended. Johns Hopkins APL: refrigerant-free future.", SAMSUNG_BLUE),
    ("AI Pro Cooking", "Internal oven camera (500\u00b0F hardened). 80 dish recognition.\nBurn detection via pixel browning rate. Personalized settings after 5 cooks.", SAMSUNG_BLUE),
    ("SmartThings Ecosystem", "Broadest integration: phones, TVs, wearables, appliances.\nMatter 1.5 (camera support first). Knox Matrix post-quantum security.", SAMSUNG_BLUE),
    ("Samsung Food + Gemini", "160K+ recipes, 40K+ ingredient recognition via phone camera.\n7-day meal plans, grocery lists. Google Gemini multimodal AI.", SAMSUNG_BLUE),
]

card_h = Inches(0.73)
card_gap = Inches(0.08)
for i, (title, desc, color) in enumerate(samsung_features):
    y = Inches(1.85) + i * (card_h + card_gap)
    # Color accent bar
    add_shape(slide1, Inches(0.5), y, Inches(0.08), card_h, color)
    # Card background
    add_shape(slide1, Inches(0.58), y, Inches(5.92), card_h, LIGHT_GRAY, BORDER_GRAY)
    add_text_box(slide1, Inches(0.7), y + Inches(0.05), Inches(5.6), Inches(0.2), title, 10, SAMSUNG_DARK, True)
    add_text_box(slide1, Inches(0.7), y + Inches(0.25), Inches(5.6), Inches(0.45), desc, 8, MED_GRAY, False)

# RIGHT: Competitor Ecosystem & Positioning
add_text_box(slide1, Inches(7), Inches(1.55), Inches(6), Inches(0.3),
             "COMPETITOR PRODUCT POSITIONING", 11, DARK_GRAY, True)

competitors = [
    ("LG", "ThinQ + Affectionate Intelligence", [
        "#1 US appliance brand (surveys)",
        "InstaView T-OLED 36\" panel, AI Fresh pre-cooling",
        "Gourmet AI: 85+ dishes (highest count)",
        "EXAONE 1.2B on-device LLM + Microsoft cloud AI",
        "Inverter Linear Compressor: best energy efficiency",
    ], LG_RED),
    ("Bosch/BSH", "Cook AI + German Engineering", [
        "#1 most trusted brand (7 years running)",
        "Cook AI: industry's most advanced agentic cooking",
        "Lambda probe O2 sensor + PerfectBake humidity",
        "Matter first-mover; open APIs",
        "CrystalDry zeolite dishwashers (unique tech)",
    ], BOSCH_RED),
    ("GE/Haier", "SmartHQ + 6-Brand Portfolio", [
        "Fastest share growth (+2 ppt, 2025)",
        "Scan-to-List: built-in barcode scanner (4M+ products)",
        "Flavorly AI: generative recipes via Google Vertex AI",
        "CookCam AI + FridgeFocus camera",
        "6 brands: Hotpoint to Monogram ($400-$12K+)",
    ], GE_BLUE),
    ("Whirlpool/KA", "Yummly + KitchenAid Relaunch", [
        "22% US share (#1 overall market)",
        "KitchenAid Intelligent Cooking Camera (KBIS 2026)",
        "Yummly: 27M+ users, largest recipe community",
        "Doneness detection, no-preheat cooking",
        "Largest product relaunch in 10+ years",
    ], WHIRLPOOL_GOLD),
    ("Miele", "20-Year Engineering Standard", [
        "Smart Food ID: 50 dishes, continuous learning",
        "M Sense Cookware: 3 temp sensors (world first)",
        "AutoDos PowerDisk: 30% detergent reduction",
        "AI Diagnostics: 80%+ self-resolution rate",
        "All 60cm dishwashers rated A for energy",
    ], MIELE_BROWN),
]

comp_h = Inches(0.95)
comp_gap = Inches(0.06)
for i, (name, tagline, bullets, color) in enumerate(competitors):
    y = Inches(1.85) + i * (comp_h + comp_gap)
    # Brand color accent
    add_shape(slide1, Inches(7), y, Inches(0.08), comp_h, color)
    add_shape(slide1, Inches(7.08), y, Inches(5.72), comp_h, LIGHT_GRAY, BORDER_GRAY)
    add_text_box(slide1, Inches(7.2), y + Inches(0.04), Inches(2), Inches(0.2), name, 10, color, True)
    add_text_box(slide1, Inches(8.8), y + Inches(0.04), Inches(3.8), Inches(0.2), tagline, 8, MED_GRAY, False)
    bullet_text = "\n".join([f"\u2022 {b}" for b in bullets])
    add_text_box(slide1, Inches(7.2), y + Inches(0.24), Inches(5.4), Inches(0.7), bullet_text, 7.5, DARK_GRAY, False)


# ============================================================
# SLIDE 2: Product & Feature Comparison by Category
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

add_shape(slide2, Inches(0), Inches(0), SLIDE_W, Inches(0.7), SAMSUNG_BLUE)
add_text_box(slide2, Inches(0.5), Inches(0.15), Inches(12), Inches(0.35),
             "Product & Feature Comparison by Category", 20, WHITE, True)

# Category comparison tables
categories = [
    {
        "name": "REFRIGERATION",
        "headers": ["Feature", "Samsung", "LG", "GE/Haier", "Bosch", "Miele"],
        "rows": [
            ["Flagship Display", "32\" FHD Touch", "T-OLED 36\"", "7\" Swivel LCD", "None", "None"],
            ["Internal Camera", "Dual-lens + Gemini", "Remote viewing", "FridgeFocus AI", "No", "FoodView (4 cam)"],
            ["Food Recognition", "1000s (Gemini)", "Inference-based", "Barcode (4M+)", "N/A", "Remote view only"],
            ["Cooling Innovation", "AI Hybrid (Peltier)", "Inverter Linear", "Standard", "Standard", "MasterCool"],
            ["Design/Custom", "Bespoke 8-14 colors", "MoodUP LED", "Cafe hardware", "Panel-ready", "Built-in"],
            ["Price (Flagship)", "$4,699", "TBD (Signature)", "$4,899", "~$3,500", "$5,000+"],
        ],
        "y": Inches(0.85),
    },
    {
        "name": "COOKING / OVENS",
        "headers": ["Feature", "Samsung", "LG", "GE/Haier", "Bosch", "Miele"],
        "rows": [
            ["Oven Camera", "Yes (500\u00b0F rated)", "Yes (Gourmet AI)", "CookCam AI", "Cook AI probe", "Smart Food ID"],
            ["Dishes Recognized", "80", "85+ (highest)", "Growing", "80 (agentic)", "50"],
            ["Autonomous Level", "Med (fridge\u2192oven)", "Med (auto-select)", "Med (precision)", "HIGH (multi-appl)", "Med (cookware)"],
            ["Unique Sensor", "Pixel browning rate", "Standard camera", "Wireless probe", "Lambda O2 probe", "M Sense (3 temp)"],
            ["Cooking Method", "Vision + Gemini AI", "Vision + LLM", "Vision + Vertex AI", "Probe + humidity", "Cookware comms"],
            ["Price (Range)", "$1,400-$3,000", "$2,000-$3,500", "$1,500-$4,000", "$2,000-$5,500+", "$3,000-$6,000+"],
        ],
        "y": Inches(2.75),
    },
    {
        "name": "AI / SOFTWARE PLATFORMS",
        "headers": ["Feature", "Samsung", "LG", "GE/Haier", "Bosch", "Whirlpool"],
        "rows": [
            ["Platform", "SmartThings", "ThinQ", "SmartHQ", "Home Connect", "Whirlpool/KA App"],
            ["LLM / GenAI", "Google Gemini", "EXAONE (on-device)", "Google Vertex AI", "Agentic (internal)", "None"],
            ["Voice AI", "Bixby + Alexa + Perplexity", "ThinQ Voice", "Alexa + Google", "Alexa+", "Alexa + Google"],
            ["Recipe Engine", "Samsung Food (160K+)", "ThinQ Recipes", "Flavorly AI", "Cook AI dynamic", "Yummly (27M users)"],
            ["Ecosystem Breadth", "Phone+TV+Watch+Home", "TV+Home", "Appliances only", "Appliances only", "Appliances only"],
            ["Matter Support", "1.5 (camera first)", "Yes", "Yes", "First-mover", "Limited"],
        ],
        "y": Inches(4.65),
    },
]

for cat in categories:
    y = cat["y"]
    # Category label
    add_shape(slide2, Inches(0.3), y, Inches(2.2), Inches(0.3), SAMSUNG_DARK)
    add_text_box(slide2, Inches(0.4), y + Inches(0.03), Inches(2), Inches(0.25), cat["name"], 10, WHITE, True)

    # Table
    table_top = y + Inches(0.35)
    cols = len(cat["headers"])
    col_widths = [Inches(1.6)] + [Inches(2.15)] * (cols - 1)
    total_w = sum([w for w in col_widths])
    row_h = Inches(0.24)

    # Header row
    x = Inches(0.3)
    brand_colors = [SAMSUNG_DARK, SAMSUNG_BLUE, LG_RED, GE_BLUE, BOSCH_RED, WHIRLPOOL_GOLD, MIELE_BROWN]
    for j, hdr in enumerate(cat["headers"]):
        w = col_widths[j]
        bg_color = brand_colors[j] if j > 0 else DARK_GRAY
        add_shape(slide2, x, table_top, w, row_h, bg_color)
        add_text_box(slide2, x + Inches(0.05), table_top + Inches(0.02), w - Inches(0.1), row_h,
                     hdr, 8, WHITE, True, PP_ALIGN.CENTER)
        x += w

    # Data rows
    for ri, row in enumerate(cat["rows"]):
        ry = table_top + (ri + 1) * row_h
        x = Inches(0.3)
        row_bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xF9, 0xFC)
        for j, cell in enumerate(row):
            w = col_widths[j]
            add_shape(slide2, x, ry, w, row_h, row_bg, BORDER_GRAY, Pt(0.5))
            font_bold = (j == 0)
            font_color = DARK_GRAY if j == 0 else BLACK
            # Highlight Samsung column
            if j == 1 and cat["name"] != "AI / SOFTWARE PLATFORMS":
                font_color = SAMSUNG_DARK
                font_bold = True
            elif j == 1:
                font_color = SAMSUNG_DARK
                font_bold = True
            add_text_box(slide2, x + Inches(0.05), ry + Inches(0.02), w - Inches(0.1), row_h,
                         cell, 7.5, font_color, font_bold, PP_ALIGN.CENTER)
            x += w

# Key insight box
add_rounded_rect(slide2, Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.65), RGBColor(0xFD, 0xF2, 0xE9))
add_text_box(slide2, Inches(0.5), Inches(6.68), Inches(12.3), Inches(0.2),
             "KEY INSIGHT", 9, RGBColor(0xE6, 0x7E, 0x22), True)
add_text_box(slide2, Inches(0.5), Inches(6.88), Inches(12.3), Inches(0.35),
             "Samsung leads in ecosystem breadth and display integration, but Bosch Cook AI sets the autonomy benchmark. "
             "GE's barcode scanner is a unique hardware innovation Samsung lacks. LG's on-device EXAONE LLM addresses cloud dependency risk Samsung has not yet solved.",
             8.5, DARK_GRAY, False)


# ============================================================
# SLIDE 3: Tech Architecture & AI Landscape
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_shape(slide3, Inches(0), Inches(0), SLIDE_W, Inches(0.7), SAMSUNG_BLUE)
add_text_box(slide3, Inches(0.5), Inches(0.15), Inches(12), Inches(0.35),
             "Technology Architecture & AI Strategy Landscape", 20, WHITE, True)

# LEFT SECTION: Architecture comparison
add_text_box(slide3, Inches(0.4), Inches(0.85), Inches(6), Inches(0.3),
             "AI ARCHITECTURE: THREE COMPETING PHILOSOPHIES", 11, SAMSUNG_DARK, True)

arch_cards = [
    ("Samsung: Distributed Edge Intelligence", SAMSUNG_BLUE, [
        "Modified Exynos SoC with dedicated NPU in each appliance",
        "L1 (Local): Real-time object detection at <50ms latency",
        "L2 (Cloud): Google Gemini for complex reasoning",
        "Advantage: Works offline, fastest response time",
        "Risk: Higher BOM cost, hard to upgrade chipset",
    ]),
    ("LG: Centralized Hub Intelligence", LG_RED, [
        "ThinQ ON hub with DQ-X AI Chipset as central brain",
        "EXAONE 1.2B-param on-device SLM for local inference",
        "Dumb node, smart hub topology via Matter/Thread",
        "Advantage: Lower appliance cost, easy hub upgrades",
        "Risk: Single point of failure, 100-300ms latency",
    ]),
    ("Bosch: Physics-First Sensor Intelligence", BOSCH_RED, [
        "Lambda probe (automotive O2 sensor) for cooking precision",
        "PerfectBake humidity sensor for moisture-curve baking",
        "AI-assisted probe thermometry, not vision-based",
        "Advantage: Superior cooking results, cloud-independent",
        "Risk: Less \"visible\" AI, lower marketing appeal",
    ]),
]

for i, (title, color, bullets) in enumerate(arch_cards):
    y = Inches(1.2) + i * Inches(1.5)
    add_shape(slide3, Inches(0.4), y, Inches(0.08), Inches(1.35), color)
    add_rounded_rect(slide3, Inches(0.48), y, Inches(5.82), Inches(1.35), LIGHT_GRAY, BORDER_GRAY)
    add_text_box(slide3, Inches(0.6), y + Inches(0.06), Inches(5.5), Inches(0.22), title, 10, color, True)
    bullet_text = "\n".join([f"\u2022 {b}" for b in bullets])
    add_text_box(slide3, Inches(0.6), y + Inches(0.3), Inches(5.5), Inches(1.0), bullet_text, 8, DARK_GRAY, False)

# Architecture comparison table
add_text_box(slide3, Inches(0.4), Inches(5.75), Inches(6), Inches(0.25),
             "HEAD-TO-HEAD: ARCHITECTURE COMPARISON", 9, SAMSUNG_DARK, True)

arch_table = [
    ["Metric", "Samsung (Edge NPU)", "LG (Central Hub)", "Bosch (Sensor)"],
    ["Vision Processing", "Superior (local)", "Cloud-dependent", "N/A (probe-based)"],
    ["Latency", "<50ms", "100-300ms", "Real-time (wired)"],
    ["Cost Structure", "High CapEx/unit", "Low CapEx, High OpEx", "Low (no NPU)"],
    ["Upgradeability", "Hard (fixed chip)", "Easy (swap hub)", "Firmware OTA"],
    ["Cloud Dependency", "Medium (Gemini)", "Medium (Microsoft)", "Low (local-first)"],
]

at_y = Inches(6.0)
at_cols = [Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.5)]
for ri, row in enumerate(arch_table):
    x = Inches(0.4)
    for ci, cell in enumerate(row):
        w = at_cols[ci]
        bg = SAMSUNG_DARK if ri == 0 else (WHITE if ri % 2 == 1 else RGBColor(0xF5, 0xF5, 0xF5))
        fc = WHITE if ri == 0 else DARK_GRAY
        fb = ri == 0 or ci == 0
        add_shape(slide3, x, at_y + ri * Inches(0.22), w, Inches(0.22), bg, BORDER_GRAY, Pt(0.5))
        add_text_box(slide3, x + Inches(0.04), at_y + ri * Inches(0.22) + Inches(0.02), w - Inches(0.08), Inches(0.2),
                     cell, 7, fc, fb, PP_ALIGN.CENTER)
        x += w

# RIGHT SECTION: Protocol & KPI
add_text_box(slide3, Inches(6.8), Inches(0.85), Inches(6), Inches(0.3),
             "COMPOSITE KPI SCORECARD (1-5 SCALE)", 11, SAMSUNG_DARK, True)

# KPI scorecard
kpi_data = [
    ("Dimension", "Wt", "Samsung", "LG", "GE", "Bosch", "Whirlpool", "Miele"),
    ("Product", "20%", "4.75", "3.85", "3.30", "3.45", "2.80", "2.80"),
    ("Technology", "20%", "5.00", "4.00", "2.40", "3.00", "2.10", "2.10"),
    ("Market", "25%", "3.70", "3.85", "3.30", "3.70", "4.10", "3.05"),
    ("Innovation", "15%", "5.00", "4.15", "2.60", "3.45", "2.50", "1.90"),
    ("Business", "20%", "4.75", "4.00", "3.55", "3.75", "2.50", "2.80"),
    ("COMPOSITE", "100%", "4.55", "3.96", "3.05", "3.48", "2.88", "2.57"),
]

kpi_cols = [Inches(0.9), Inches(0.4), Inches(0.8), Inches(0.6), Inches(0.6), Inches(0.65), Inches(0.9), Inches(0.6)]
kpi_brand_colors = [DARK_GRAY, DARK_GRAY, SAMSUNG_BLUE, LG_RED, GE_BLUE, BOSCH_RED, WHIRLPOOL_GOLD, MIELE_BROWN]

for ri, row in enumerate(kpi_data):
    x = Inches(6.8)
    for ci, cell in enumerate(row):
        w = kpi_cols[ci]
        is_header = (ri == 0)
        is_composite = (ri == len(kpi_data) - 1)
        if is_header:
            bg = kpi_brand_colors[ci] if ci >= 2 else DARK_GRAY
        elif is_composite:
            bg = RGBColor(0xE8, 0xEB, 0xF7) if ci >= 2 else RGBColor(0xE0, 0xE0, 0xE0)
        else:
            bg = WHITE if ri % 2 == 1 else RGBColor(0xF8, 0xF8, 0xF8)
        fc = WHITE if is_header else (SAMSUNG_DARK if (ci == 2 and is_composite) else DARK_GRAY)
        fb = is_header or ci == 0 or is_composite
        add_shape(slide3, x, Inches(1.15) + ri * Inches(0.26), w, Inches(0.26), bg, BORDER_GRAY, Pt(0.5))
        add_text_box(slide3, x + Inches(0.03), Inches(1.17) + ri * Inches(0.26), w - Inches(0.06), Inches(0.22),
                     cell, 7.5, fc, fb, PP_ALIGN.CENTER)
        x += w

# Protocol / Smart Home comparison
add_text_box(slide3, Inches(6.8), Inches(3.2), Inches(6), Inches(0.25),
             "SMART HOME PROTOCOL SUPPORT", 10, SAMSUNG_DARK, True)

protocol_data = [
    ["Protocol", "Samsung", "LG", "Bosch", "GE", "Miele"],
    ["Matter", "1.5 (camera)", "Yes", "First-mover", "Yes", "Limited"],
    ["Thread", "Yes", "Yes", "In progress", "In progress", "Limited"],
    ["Voice Native", "Bixby+Alexa", "ThinQ", "Alexa+", "SmartHQ", "App only"],
    ["LLM Partner", "Google Gemini", "Microsoft", "Internal", "Google Cloud", "None"],
    ["Security", "Knox Matrix", "LG Shield", "GDPR std", "Google Cloud", "Standard"],
]

pt_cols = [Inches(1.0), Inches(1.1), Inches(0.8), Inches(1.0), Inches(0.9), Inches(0.8)]
pt_brand_colors = [DARK_GRAY, SAMSUNG_BLUE, LG_RED, BOSCH_RED, GE_BLUE, MIELE_BROWN]
for ri, row in enumerate(protocol_data):
    x = Inches(6.8)
    for ci, cell in enumerate(row):
        w = pt_cols[ci]
        bg = pt_brand_colors[ci] if ri == 0 else (WHITE if ri % 2 == 1 else RGBColor(0xF8, 0xF8, 0xF8))
        fc = WHITE if ri == 0 else DARK_GRAY
        fb = ri == 0 or ci == 0
        add_shape(slide3, x, Inches(3.45) + ri * Inches(0.24), w, Inches(0.24), bg, BORDER_GRAY, Pt(0.5))
        add_text_box(slide3, x + Inches(0.03), Inches(3.47) + ri * Inches(0.24), w - Inches(0.06), Inches(0.2),
                     cell, 7.5, fc, fb, PP_ALIGN.CENTER)
        x += w

# CES 2026 Awards & Highlights
add_text_box(slide3, Inches(6.8), Inches(5.0), Inches(6), Inches(0.25),
             "CES 2026 KITCHEN AI HIGHLIGHTS", 10, SAMSUNG_DARK, True)

ces_items = [
    ("Samsung", "CES Innovation Award. First appliance with Google Gemini.\nAuto Open Door voice-controlled mechanism.", SAMSUNG_BLUE),
    ("LG", "\"Affectionate Intelligence\" theme. LG SIGNATURE with\nbuilt-in LLM. ThinQ Pro for B2B.", LG_RED),
    ("Bosch", "Reviewed CES Award. Cook AI agentic cooking.\nAlexa+ espresso machine.", BOSCH_RED),
    ("GE", "Smart Company of Year (9x). Kitchen Assistant\nwith barcode scanner. Instacart integration.", GE_BLUE),
    ("Industry", "Shift: Smart \u2192 Autonomous. 239 smart kitchen\nstartups globally. LLMs standard in premium.", DARK_GRAY),
]

for i, (brand, desc, color) in enumerate(ces_items):
    y = Inches(5.3) + i * Inches(0.42)
    add_shape(slide3, Inches(6.8), y, Inches(0.07), Inches(0.37), color)
    add_text_box(slide3, Inches(6.95), y + Inches(0.01), Inches(0.9), Inches(0.2), brand, 8, color, True)
    add_text_box(slide3, Inches(7.9), y + Inches(0.01), Inches(4.7), Inches(0.35), desc, 7, DARK_GRAY, False)


# ============================================================
# SLIDE 4: Competitive Gaps & Strategic Priorities
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_shape(slide4, Inches(0), Inches(0), SLIDE_W, Inches(0.7), SAMSUNG_BLUE)
add_text_box(slide4, Inches(0.5), Inches(0.15), Inches(12), Inches(0.35),
             "Competitive Gaps & Strategic Priorities (2026-2030)", 20, WHITE, True)

# LEFT: Competitive Gaps
add_text_box(slide4, Inches(0.4), Inches(0.85), Inches(6), Inches(0.3),
             "KEY COMPETITIVE GAPS FOR SAMSUNG", 11, RED_ALERT, True)

gaps = [
    ("HIGH", "Autonomous Cooking AI", "Bosch",
     "Bosch Cook AI is industry's most advanced agentic system.\nOrchestrates multiple appliances, adapts to ingredients,\nuses lambda O2 sensor. Samsung is vision-focused, not agentic.", BOSCH_RED, RED_ALERT),
    ("HIGH", "Brand Trust Gap", "Bosch",
     "Samsung ranks 4th in kitchen trust. Bosch #1 for 7 years.\nConnected appliances: 87 problems/100 units vs 63 non-connected.\nReliability perception hasn't caught up with innovation.", BOSCH_RED, RED_ALERT),
    ("HIGH", "GE Market Momentum", "GE/Haier",
     "Fastest share growth (+2 ppt, 2025). 9x Smart Company of Year.\nScan-to-List barcode scanner is unique innovation.\nInstacart integration = practical kitchen commerce.", GE_BLUE, RED_ALERT),
    ("MED-HIGH", "On-Device AI / Cloud Risk", "LG",
     "Samsung's Gemini dependency creates 15-20 year lifespan risk.\nLG's EXAONE runs locally. Bosch is sensor-local.\n47% of smart kitchen users worry about data breaches.", LG_RED, AMBER),
    ("MED-HIGH", "KitchenAid Smart Relaunch", "Whirlpool",
     "Largest product relaunch in 10+ years at KBIS 2026.\nIntelligent Cooking Camera directly competes with Samsung.\nYummly's 27M users = massive recipe data moat.", WHIRLPOOL_GOLD, AMBER),
    ("MEDIUM", "Luxury Segment (Dacor)", "Multiple",
     "Dacor hasn't matched Thermador, Sub-Zero/Wolf, or Miele prestige.\nSub-Zero dual refrigeration is technically superior.\nMiele's 20-year engineering standard is unmatched.", MIELE_BROWN, AMBER),
    ("MEDIUM", "Chinese Brand Disruption", "Midea",
     "Midea overseas revenue +17.7% YoY. New US leadership team.\nMidea-Hisense pact pools AI + manufacturing.\nAggressive pricing could pressure Samsung's mid-range.", DARK_GRAY, AMBER),
]

for i, (severity, title, competitor, desc, brand_color, sev_color) in enumerate(gaps):
    y = Inches(1.2) + i * Inches(0.85)
    # Severity badge
    sev_w = Inches(0.75)
    add_rounded_rect(slide4, Inches(0.4), y + Inches(0.04), sev_w, Inches(0.2), sev_color)
    add_text_box(slide4, Inches(0.4), y + Inches(0.04), sev_w, Inches(0.2), severity, 7, WHITE, True, PP_ALIGN.CENTER)
    # Title + competitor
    add_text_box(slide4, Inches(1.2), y + Inches(0.02), Inches(3.5), Inches(0.22), title, 10, DARK_GRAY, True)
    add_text_box(slide4, Inches(4.5), y + Inches(0.04), Inches(1.5), Inches(0.2), f"vs {competitor}", 8, brand_color, True)
    # Description
    add_text_box(slide4, Inches(1.2), y + Inches(0.24), Inches(5), Inches(0.55), desc, 7.5, MED_GRAY, False)

# RIGHT: Strategic Priorities
add_text_box(slide4, Inches(6.8), Inches(0.85), Inches(6), Inches(0.3),
             "7 STRATEGIC PRIORITIES", 11, SAMSUNG_DARK, True)

priorities = [
    ("P1", "Close Reliability/Trust Gap", "CRITICAL",
     "Launch Samsung Reliability Guarantee. Target 4.5+ app rating.\nAddress 87/100 connected problem rate. Publish reliability data.", RED_ALERT),
    ("P2", "Advance Autonomous Cooking AI", "HIGH",
     "Develop agentic cooking AI matching Bosch Cook AI.\nLeverage Gemini multimodal for ingredient-adaptive cooking.\nFridge\u2192Recipe\u2192Oven fully automated orchestration.", RED_ALERT),
    ("P3", "Defend Smart Kitchen Share vs GE", "HIGH",
     "Match barcode scanner hardware innovation.\nDeepen grocery delivery integration from Family Hub.\nLeverage phone+TV+watch ecosystem GE can't match.", RED_ALERT),
    ("P4", "Strengthen Matter/Interop Strategy", "MED-HIGH",
     "Position SmartThings as best Matter hub/controller.\nAccelerate Matter across all price tiers (not just premium).\nEnsure seamless Google Home + Apple HomeKit via Matter.", AMBER),
    ("P5", "Address Cloud Dependency Risk", "MEDIUM",
     "Develop hybrid AI: core functions on-device, cloud enhances.\nEstablish long-term Gemini SLA for appliance support.\nPlan graceful degradation if cloud services change.", AMBER),
    ("P6", "Elevate Dacor in Luxury", "MEDIUM",
     "Infuse Samsung AI + SmartThings into Dacor products.\nCreate 'Dacor x Samsung AI' luxury smart suite.\nDifferentiate as AI luxury vs Sub-Zero/Wolf's tech-avoidance.", AMBER),
    ("P7", "Counter Chinese Brand Entry", "MEDIUM",
     "Protect mid-range with compelling smart features.\nAccelerate Bespoke entry-tier expansion.\nLeverage ecosystem depth as competitive moat vs Midea.", AMBER),
]

for i, (num, title, severity, desc, sev_color) in enumerate(priorities):
    y = Inches(1.15) + i * Inches(0.82)
    # Number circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), y + Inches(0.05), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SAMSUNG_BLUE
    circle.line.fill.background()
    add_text_box(slide4, Inches(6.8), y + Inches(0.08), Inches(0.3), Inches(0.25), num, 8, WHITE, True, PP_ALIGN.CENTER)
    # Title
    add_text_box(slide4, Inches(7.15), y + Inches(0.05), Inches(3.5), Inches(0.22), title, 10, DARK_GRAY, True)
    # Severity badge
    add_rounded_rect(slide4, Inches(10.5), y + Inches(0.06), Inches(0.8), Inches(0.2), sev_color)
    add_text_box(slide4, Inches(10.5), y + Inches(0.06), Inches(0.8), Inches(0.2), severity, 7, WHITE, True, PP_ALIGN.CENTER)
    # Description
    add_text_box(slide4, Inches(7.15), y + Inches(0.28), Inches(5.5), Inches(0.5), desc, 7.5, MED_GRAY, False)

# Bottom outlook bar
add_shape(slide4, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), SAMSUNG_DARK)
add_text_box(slide4, Inches(0.5), Inches(7.1), Inches(3), Inches(0.2),
             "2026-2030 OUTLOOK", 10, WHITE, True)
outlook_text = (
    "2026: Smart kitchen tipping point \u2022 2027: Autonomous cooking goes commercial + Samsung proprietary GPU AP \u2022 "
    "2028: 25% smart penetration, edge AI 80%+ inference \u2022 2029: 30.8% global smart homes \u2022 "
    "2030: $40-60B smart kitchen market, refrigerant-free Peltier cooling"
)
add_text_box(slide4, Inches(3.5), Inches(7.1), Inches(9.3), Inches(0.35), outlook_text, 7.5, RGBColor(0xB0, 0xBE, 0xE8), False)


# ============================================================
# SLIDE 5: Appendix - Data Sources
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

add_shape(slide5, Inches(0), Inches(0), SLIDE_W, Inches(0.7), SAMSUNG_DARK)
add_text_box(slide5, Inches(0.5), Inches(0.15), Inches(12), Inches(0.35),
             "Appendix: Data Sources & References", 20, WHITE, True)
add_text_box(slide5, Inches(0.5), Inches(0.5), Inches(10), Inches(0.2),
             "100+ primary sources across market research firms, company newsrooms, trade publications, and technology analysis", 9, RGBColor(0xB0, 0xBE, 0xE8), False)

# Source categories
source_sections = [
    ("MARKET DATA & SIZING", SAMSUNG_BLUE, [
        ("$311B global kitchen market, CAGR 4.2-4.8%", "Mordor Intelligence, Coherent Market Insights, IMARC Group (2025-2026)"),
        ("Smart kitchen $23.7B (2025) \u2192 $40-60B (2030)", "Research and Markets, Statista, Polaris Market Research"),
        ("US market ~$77-80B; smart kitchen $6.9B \u2192 $17.64B by 2033", "GlobeNewsWire US Smart Kitchen Forecast 2025-2033"),
        ("Smart household penetration 13% \u2192 30.8% by 2029", "Statista Digital Market Insights, American Home Shield Survey"),
    ]),
    ("MARKET SHARE & BRAND DATA", RGBColor(0x00, 0x7A, 0x33), [
        ("Samsung 18% US / 21.7% smart kitchen share", "OpenBrand US Major Appliance Market Share Q1-Q2 2025"),
        ("Whirlpool 22%, LG 17%, GE 14%, Bosch 7%", "OpenBrand Q2 2025, Fred's Appliance Academy Industry Results"),
        ("GE +2 ppt share gain (fastest growth, 2025)", "OpenBrand Q2 2025 Market Share Trends & Rankings"),
        ("Bosch #1 most trusted kitchen brand (7 years)", "Lifestory Research 2025 Most Trusted Study"),
        ("87 problems/100 connected units; satisfaction -11 pts", "J.D. Power 2025 U.S. Appliance Satisfaction Study"),
    ]),
    ("SAMSUNG PRODUCT & AI DATA", SAMSUNG_BLUE, [
        ("Bespoke AI Family Hub+, Gemini, AI Hybrid Cooling", "Samsung Global Newsroom (news.samsung.com), CES 2026 press releases"),
        ("AI Vision Inside, 80 dish recognition, SmartThings 425M+", "Samsung Bespoke AI Kitchen Appliances press materials"),
        ("Samsung Food 160K+ recipes, Matter 1.5 camera-first", "Samsung CES 2026 announcements, SmartThings Blog"),
        ("Knox Matrix, Mach-1 AI Chip, 9,304 US patents (2024)", "Samsung SEC filings, GreyB Patent Analytics"),
    ]),
    ("COMPETITOR PRODUCT & AI DATA", MED_GRAY, [
        ("LG: Gourmet AI 85+ dishes, EXAONE 1.2B LLM, ThinQ ON", "LG Newsroom CES 2026, LG SIGNATURE press releases"),
        ("Bosch: Cook AI agentic system, lambda probe, CrystalDry", "BSH CES 2026 (BusinessWire), Reviewed CES coverage"),
        ("GE: Kitchen Assistant, barcode scanner 4M+, Flavorly AI", "GE Appliances Press Room, Google Cloud AI partnership release"),
        ("Whirlpool: KitchenAid Camera AI, Yummly 27M+, KBIS 2026", "Whirlpool Corp KBIS 2026 (PRNewswire), KitchenAid press"),
        ("Miele: Smart Food ID, M Sense cookware, 20-year standard", "Miele SBID Kitchen Trends 2026, mieleusa.com"),
    ]),
    ("TECHNOLOGY & ARCHITECTURE", RGBColor(0x8E, 0x44, 0xAD), [
        ("Edge NPU vs Central Hub vs Sensor-first architectures", "Gemini Market Intelligence Report (custom research, 2026)"),
        ("Matter protocol versions 1.2-1.5, 750+ certified products", "Connectivity Standards Alliance (CSA), BSH/Samsung press"),
        ("AI patent growth 800%+ since 2017", "GreyB Patent Analytics, Samsung/LG/BSH patent filings"),
        ("Cloud dependency risk, 47% data breach concern", "American Home Shield Smart Home Survey, industry analysis"),
    ]),
    ("CES 2026 & INDUSTRY EVENTS", RGBColor(0xE6, 0x7E, 0x22), [
        ("CES 2026 Innovation Awards, product announcements", "CES official, Reviewed, Tom's Guide, Engadget, The Kitchn"),
        ("KBIS 2026 Best of Show, KitchenAid/Whirlpool launches", "KBIS 2026 official, PRNewswire KBIS coverage"),
        ("GE 9x Smart Appliance Company of Year", "GE Appliances Press Room (2026 award announcement)"),
        ("239 smart kitchen startups, 96 in US, 35 Series A+", "Tracxn Smart Kitchen Appliances Report"),
    ]),
]

col_w = Inches(6.2)
row_h_header = Inches(0.25)
row_h_item = Inches(0.28)
gap_between_sections = Inches(0.05)

# Lay out in 2 columns
y_positions = [Inches(0.85), Inches(0.85)]  # track Y for each column

for si, (section_title, section_color, items) in enumerate(source_sections):
    col = 0 if si < 3 else 1
    x = Inches(0.3) if col == 0 else Inches(6.8)
    y = y_positions[col]

    # Section header
    add_shape(slide5, x, y, col_w, row_h_header, section_color)
    add_text_box(slide5, x + Inches(0.1), y + Inches(0.02), col_w - Inches(0.2), row_h_header,
                 section_title, 8, WHITE, True)
    y += row_h_header + Inches(0.03)

    for data_point, source in items:
        # Data point
        add_text_box(slide5, x + Inches(0.1), y, col_w - Inches(0.2), Inches(0.14),
                     f"\u2022 {data_point}", 7, DARK_GRAY, True)
        # Source
        add_text_box(slide5, x + Inches(0.2), y + Inches(0.13), col_w - Inches(0.3), Inches(0.14),
                     f"Source: {source}", 6.5, MED_GRAY, False)
        y += row_h_item

    y += gap_between_sections
    y_positions[col] = y

# Footer note
add_shape(slide5, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), RGBColor(0xF0, 0xF2, 0xFA))
add_text_box(slide5, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             "Report compiled March 2026 based on publicly available information. All market figures are directional estimates. "
             "KPI scoring uses 1-5 scale based on public data. Full source list (60+ URLs) available in companion research document.",
             7.5, MED_GRAY, False)


# Save
output_path = "/home/chachapranu/ashu/git/product-research/claude/Samsung_Kitchen_Comprehensive_2026.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
