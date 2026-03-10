#!/usr/bin/env python3
"""
Samsung vs Competition: AI Kitchen Features - 4-Slide Executive Presentation
Target audience: VP & Director level executives
Focus: Competitive gaps and strategic landscape
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)
SAMSUNG_LIGHT = RGBColor(0x1E, 0x3F, 0xC2)
SAMSUNG_DARK = RGBColor(0x0D, 0x1A, 0x6B)
LG_RED = RGBColor(0xA5, 0x00, 0x34)
BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
GE_BLUE = RGBColor(0x00, 0x4B, 0x8D)
WHIRLPOOL_GOLD = RGBColor(0xC8, 0x96, 0x23)
MIELE_BROWN = RGBColor(0x6E, 0x1D, 0x24)
HAIER_BLUE = RGBColor(0x00, 0x7D, 0xB8)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF5)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
GREEN = RGBColor(0x0F, 0x7B, 0x3F)
RED_ALERT = RGBColor(0xD4, 0x2B, 0x2B)
AMBER = RGBColor(0xE6, 0x8A, 0x00)

# Star ratings
STAR_GOLD = RGBColor(0xF5, 0xA6, 0x23)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(text_frame, items, font_size=11, color=BLACK, bold_prefix=None, spacing=Pt(4)):
    """Add bullet items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.space_after = spacing
        p.space_before = Pt(1)

        if isinstance(item, tuple):
            # (bold_part, normal_part)
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = "Calibri"
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Calibri"


def add_icon_card(slide, left, top, width, height, icon_text, title, subtitle, icon_color, bg_color=WHITE):
    """Create a card with colored icon circle, title, and subtitle."""
    card = add_rounded_rect(slide, left, top, width, height, bg_color)

    # Icon circle
    circle_size = Inches(0.45)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.18), circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # Title
    add_text_box(slide, left + Inches(0.75), top + Inches(0.12),
                 width - Inches(0.9), Inches(0.35), title,
                 font_size=12, color=DARK_GRAY, bold=True)

    # Subtitle
    add_text_box(slide, left + Inches(0.75), top + Inches(0.42),
                 width - Inches(0.9), height - Inches(0.5), subtitle,
                 font_size=9, color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Samsung AI Kitchen Position
# ═══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, WHITE)

# Top banner
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Inches(1.5), SAMSUNG_BLUE)

# Title
add_text_box(slide1, Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
             "Samsung AI Kitchen: Competitive Landscape 2026",
             font_size=26, color=WHITE, bold=True)

# Subtitle
add_text_box(slide1, Inches(0.6), Inches(0.78), Inches(10), Inches(0.5),
             "AI Feature Deep Dive  |  Samsung vs LG, Bosch, GE/Haier, Whirlpool/KitchenAid, Miele",
             font_size=13, color=RGBColor(0xB0, 0xC4, 0xEE))

# Market context bar
add_shape(slide1, Inches(0), Inches(1.5), SLIDE_W, Inches(0.55), RGBColor(0xE8, 0xED, 0xF5))
add_text_box(slide1, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4),
             "US Smart Kitchen Market: $6.9B (2025) → $17.6B by 2033 (12.5% CAGR)    |    Samsung Share: ~21.7%    |    Industry Shift: 'Smart' → 'Autonomous' Kitchens",
             font_size=11, color=SAMSUNG_DARK, bold=True)

# Samsung Strengths - Left column header
add_text_box(slide1, Inches(0.5), Inches(2.25), Inches(6), Inches(0.4),
             "SAMSUNG'S AI STRENGTHS",
             font_size=14, color=GREEN, bold=True)

# Left accent bar
add_shape(slide1, Inches(0.5), Inches(2.65), Inches(5.8), Inches(0.04), GREEN)

# Samsung strength cards
strengths = [
    ("S", "SmartThings Ecosystem", "410M+ users | phones, TVs, wearables, appliances | broadest in industry", SAMSUNG_BLUE),
    ("G", "Google Gemini Vision AI", "First Gemini in a fridge | multimodal food recognition | 37+ foods expanding", SAMSUNG_LIGHT),
    ("K", "Knox Matrix Security", "Hardware chip + blockchain + post-quantum encryption | industry-leading", RGBColor(0x2E, 0x7D, 0x32)),
    ("H", "AI Hybrid Cooling", "Peltier + compressor dual system | Johns Hopkins APL partnership | unique", RGBColor(0x00, 0x79, 0x6B)),
    ("B", "Bixby Voice AI", "Voice ID multi-user | Perplexity AI (2026) | built-in mic/speaker", RGBColor(0x4A, 0x14, 0x8C)),
    ("D", "Display Leadership", "Up to 32\" AI Home touchscreen | 3D Map View | smart home command center", RGBColor(0x0D, 0x47, 0xA1)),
]

y_pos = Inches(2.85)
for i, (icon, title, sub, clr) in enumerate(strengths):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.5) + col * Inches(3.1)
    y = y_pos + row * Inches(0.85)
    add_icon_card(slide1, x, y, Inches(2.9), Inches(0.75), icon, title, sub, clr, LIGHT_GRAY)

# Competitor landscape - Right column header
add_text_box(slide1, Inches(7), Inches(2.25), Inches(6), Inches(0.4),
             "COMPETITOR AI PHILOSOPHIES",
             font_size=14, color=RED_ALERT, bold=True)

# Right accent bar
add_shape(slide1, Inches(7), Inches(2.65), Inches(5.8), Inches(0.04), RED_ALERT)

competitors = [
    ("LG", "Affectionate Intelligence", "On-device EXAONE LLM | 85+ dish Gourmet AI\n36\" T-OLED display | 170+ brand ecosystem", LG_RED),
    ("Bosch", "Agentic AI Cooking", "Cook AI multi-appliance orchestration\nLambda probe sensor fusion | Matter-first", BOSCH_RED),
    ("GE/Haier", "Practical AI at Every Tier", "Edge AI privacy (images on-device) | Flavorly AI\nBarcode scanner | Instacart integration", GE_BLUE),
    ("Whirlpool", "Camera AI + Yummly", "KitchenAid Intelligent Cooking Camera\n27M+ Yummly users | KBIS 2026 relaunch", WHIRLPOOL_GOLD),
    ("Miele", "Precision-Engineering AI", "M Sense smart cookware (world first)\n80%+ AI diagnostic self-resolution", MIELE_BROWN),
]

y_pos2 = Inches(2.85)
for i, (name, philosophy, detail, clr) in enumerate(competitors):
    y = y_pos2 + i * Inches(0.85)
    # Color bar
    add_shape(slide1, Inches(7), y, Inches(0.08), Inches(0.75), clr)

    # Brand name
    add_text_box(slide1, Inches(7.2), y + Inches(0.02), Inches(1.2), Inches(0.3),
                 name, font_size=11, color=clr, bold=True)

    # Philosophy
    add_text_box(slide1, Inches(8.3), y + Inches(0.02), Inches(2.5), Inches(0.3),
                 philosophy, font_size=10, color=DARK_GRAY, bold=True)

    # Detail
    add_text_box(slide1, Inches(7.2), y + Inches(0.3), Inches(5.5), Inches(0.45),
                 detail, font_size=9, color=MED_GRAY)

# Date/source footer
add_text_box(slide1, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
             "Source: CES 2025/2026, IFA 2025, KBIS 2026, manufacturer press releases, industry publications  |  March 2026",
             font_size=8, color=MED_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Category-by-Category Comparison (Data-Heavy)
# ═══════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# Header bar
add_shape(slide2, Inches(0), Inches(0), SLIDE_W, Inches(0.85), SAMSUNG_BLUE)
add_text_box(slide2, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
             "AI Feature Comparison: Samsung vs Competition",
             font_size=22, color=WHITE, bold=True)

# Table setup
col_headers = ["AI CATEGORY", "SAMSUNG", "LG", "BOSCH", "GE / HAIER", "WHIRLPOOL", "MIELE"]
col_widths = [Inches(1.8), Inches(1.95), Inches(1.85), Inches(1.85), Inches(1.95), Inches(1.75), Inches(1.55)]
col_colors = [DARK_GRAY, SAMSUNG_BLUE, LG_RED, BOSCH_RED, GE_BLUE, WHIRLPOOL_GOLD, MIELE_BROWN]

# Scoring data - using filled/empty circles
# ★★★★★ = 5, etc. Using text-based stars
rows_data = [
    ("Fridge Vision AI", "★★★★★", "★★★★", "—", "★★★★", "★", "★★★"),
    ("Oven Cooking AI", "★★★★", "★★★★★", "★★★★", "★★★", "★★★", "★★★"),
    ("Autonomous Cook", "★★★", "★★★★", "★★★★★", "★★★", "★★", "★★★"),
    ("Voice / Conv. AI", "★★★★★", "★★★★", "★★★", "★★", "★★", "★★"),
    ("Energy Optim.", "★★★★★", "★★★★", "★★", "★★★★", "★★", "★★"),
    ("Recipe AI Platform", "★★★★", "★★★", "★★★★", "★★★★★", "★★★", "★★★"),
    ("Maintenance AI", "★★★★", "★★★★", "★★★", "★★★", "★★", "★★★★★"),
    ("Sensor Innovation", "★★★★", "★★★", "★★★★★", "★★★", "★★", "★★★★★"),
    ("Ecosystem Breadth", "★★★★★", "★★★★", "★★★★", "★★★★", "★★★", "★★"),
    ("Edge AI / Privacy", "★★★", "★★★★", "★★", "★★★★★", "★★★", "★★★★"),
    ("Security", "★★★★★", "★★★", "★★★", "★★★★", "★★", "★★★"),
]

# Draw table
table_left = Inches(0.25)
table_top = Inches(1.05)
row_h = Inches(0.47)
header_h = Inches(0.45)

# Header row background
x_pos = table_left
for j, (header, w, clr) in enumerate(zip(col_headers, col_widths, col_colors)):
    add_shape(slide2, x_pos, table_top, w, header_h, clr)
    add_text_box(slide2, x_pos, table_top + Inches(0.05), w, header_h,
                 header, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    x_pos += w

# Data rows
for i, row in enumerate(rows_data):
    y = table_top + header_h + i * row_h
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE

    x_pos = table_left
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        add_shape(slide2, x_pos, y, w, row_h, bg,
                  line_color=BORDER_GRAY, line_width=0.5)

        font_sz = 9 if j == 0 else 10
        clr = DARK_GRAY if j == 0 else col_colors[j] if "★★★★★" in cell else DARK_GRAY
        bld = True if j == 0 or "★★★★★" in cell else False
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER

        text_left = x_pos + (Inches(0.1) if j == 0 else Inches(0))
        add_text_box(slide2, text_left, y + Inches(0.07), w - Inches(0.1), row_h,
                     cell, font_size=font_sz, color=clr, bold=bld, alignment=align)
        x_pos += w

# Legend & Key Insight
legend_y = table_top + header_h + len(rows_data) * row_h + Inches(0.15)

add_text_box(slide2, Inches(0.4), legend_y, Inches(5), Inches(0.3),
             "★★★★★ = Category Leader    ★★★★ = Strong    ★★★ = Competitive    ★★ = Developing    — = Not Available",
             font_size=8, color=MED_GRAY)

# Key insight box
insight_y = legend_y + Inches(0.3)
add_shape(slide2, Inches(0.25), insight_y, Inches(12.8), Inches(0.9),
          RGBColor(0xFE, 0xF3, 0xE2), line_color=AMBER, line_width=1.5)

insight_box = slide2.shapes.add_textbox(Inches(0.5), insight_y + Inches(0.08), Inches(12.3), Inches(0.75))
tf = insight_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "KEY INSIGHT: "
run1.font.size = Pt(10)
run1.font.color.rgb = AMBER
run1.font.bold = True
run1.font.name = "Calibri"
run2 = p.add_run()
run2.text = ("Samsung leads in 4 of 11 categories (Fridge Vision, Voice AI, Energy, Ecosystem, Security) but faces category-leading competition "
             "from Bosch (Autonomous Cooking, Sensor Innovation), Miele (Maintenance AI, Sensor Innovation), GE (Recipe AI, Edge Privacy), "
             "and LG (Oven AI, On-Device LLM). No single brand dominates all categories.")
run2.font.size = Pt(9.5)
run2.font.color.rgb = DARK_GRAY
run2.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: Where Competitors Excel (Competitive Gaps)
# ═══════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

# Header bar
add_shape(slide3, Inches(0), Inches(0), SLIDE_W, Inches(0.85), RGBColor(0x8B, 0x0D, 0x0D))
add_text_box(slide3, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
             "Competitive Gaps: Where Samsung Must Respond",
             font_size=22, color=WHITE, bold=True)

# Threat level indicator
add_text_box(slide3, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.4),
             "● HIGH    ● MEDIUM-HIGH    ● MEDIUM",
             font_size=9, color=RGBColor(0xFF, 0xCC, 0xCC))

# Gap cards - 2 rows x 3 columns
gaps = [
    {
        "brand": "BOSCH",
        "color": BOSCH_RED,
        "threat": "HIGH",
        "threat_color": RED_ALERT,
        "title": "Agentic AI Cooking",
        "gap": "Cook AI orchestrates cooktop + oven + probe simultaneously with minimal user input. Samsung's AI Pro Cooking is single-appliance only.",
        "detail": "Cook AI: photo ingredients → specify outcome → AI coordinates multiple appliances. Can cook multiple steaks to different doneness levels. Uses lambda probe (automotive O₂ sensor) + PerfectBake moisture sensor for triple-sensor feedback.",
        "impact": "Samsung lacks multi-appliance cooking orchestration -- a generation behind in autonomous cooking."
    },
    {
        "brand": "GE / HAIER",
        "color": GE_BLUE,
        "threat": "MED-HIGH",
        "threat_color": AMBER,
        "title": "Edge AI Privacy",
        "gap": "CookCam AI: \"Images stay on the device, not in the cloud.\" Samsung's Gemini shift sends food images to Google Cloud.",
        "detail": "GE splits architecture: visual data on-device (privacy), generative features in cloud (capability). Haier HomeGPT is Level 4 certified edge LLM -- full offline AI. Samsung moving opposite direction with cloud dependency.",
        "impact": "Growing privacy awareness + 15-20yr appliance lifespan vs. cloud API uncertainty."
    },
    {
        "brand": "GE / HAIER",
        "color": GE_BLUE,
        "threat": "MEDIUM",
        "threat_color": AMBER,
        "title": "Generative Recipe AI",
        "gap": "Flavorly AI (Google Gemini/Vertex): photo-to-recipe reverse engineering, handwritten recipe digitization, Instacart integration.",
        "detail": "Features Samsung Food doesn't offer: photograph a restaurant dish → AI recreates recipe. Photo family recipes → digital conversion. Recipe-to-Instacart delivery in 30 min. Available without buying any GE appliance.",
        "impact": "GE's recipe platform is more capable and accessible than Samsung Food."
    },
    {
        "brand": "MIELE",
        "color": MIELE_BROWN,
        "threat": "MEDIUM",
        "threat_color": AMBER,
        "title": "Smart Cookware Sensors",
        "gap": "M Sense: world's first cookware with 3 embedded temp sensors + DirectTouch controls. Samsung has no smart cookware strategy.",
        "detail": "Cookware communicates with KM 8000 induction hob via multi-patented tech. Auto-adjusts power. Prevents overcooking/burning. Ships April 2026. Likely to spawn competitor products industry-wide.",
        "impact": "New category Samsung doesn't compete in. Sensors in cookware, not just appliances."
    },
    {
        "brand": "LG",
        "color": LG_RED,
        "threat": "MEDIUM",
        "threat_color": AMBER,
        "title": "On-Device LLM in Appliances",
        "gap": "EXAONE 4.0 (1.2B params) runs locally on SIGNATURE fridge. Offers privacy + offline AI. Samsung's Bixby/Gemini is cloud-dependent.",
        "detail": "LG is first with proprietary LLM in kitchen appliance. EXAONE 4.5 (MWC 2026) adds vision-language capability. K-EXAONE global expansion with 200MW AI data center (120K GPUs). Asia first, then US.",
        "impact": "On-device LLMs will proliferate as consumer privacy awareness grows."
    },
    {
        "brand": "BOSCH / BSH",
        "color": BOSCH_RED,
        "threat": "MEDIUM",
        "threat_color": AMBER,
        "title": "Matter-First Ecosystem",
        "gap": "BSH was first globally to ship Matter-enabled appliance. Open API with 80+ features. May erode Samsung's SmartThings lock-in.",
        "detail": "As Matter adoption accelerates (750+ certified products), proprietary ecosystems lose value. BSH's open approach + 5 brands (Bosch, Siemens, Gaggenau, Neff, Thermador) creates broad interoperable network.",
        "impact": "SmartThings advantage may diminish as Matter standardizes cross-brand control."
    },
]

card_w = Inches(4.05)
card_h = Inches(2.85)
start_x = Inches(0.3)
start_y = Inches(1.1)
gap_x = Inches(0.15)
gap_y = Inches(0.15)

for i, g in enumerate(gaps):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    # Card background
    add_rounded_rect(slide3, x, y, card_w, card_h, LIGHT_GRAY)

    # Brand color top bar
    add_shape(slide3, x, y, card_w, Inches(0.06), g["color"])

    # Brand name + threat level
    add_text_box(slide3, x + Inches(0.15), y + Inches(0.1), Inches(2), Inches(0.25),
                 g["brand"], font_size=9, color=g["color"], bold=True)

    # Threat indicator
    threat_text = f"● {g['threat']}"
    add_text_box(slide3, x + card_w - Inches(1.3), y + Inches(0.1), Inches(1.15), Inches(0.25),
                 threat_text, font_size=8, color=g["threat_color"], bold=True, alignment=PP_ALIGN.RIGHT)

    # Gap title
    add_text_box(slide3, x + Inches(0.15), y + Inches(0.35), card_w - Inches(0.3), Inches(0.3),
                 g["title"], font_size=12, color=DARK_GRAY, bold=True)

    # Gap description
    add_text_box(slide3, x + Inches(0.15), y + Inches(0.65), card_w - Inches(0.3), Inches(0.55),
                 g["gap"], font_size=8.5, color=MED_GRAY)

    # Detail
    add_text_box(slide3, x + Inches(0.15), y + Inches(1.2), card_w - Inches(0.3), Inches(0.95),
                 g["detail"], font_size=8, color=MED_GRAY)

    # Impact bar
    impact_y = y + card_h - Inches(0.48)
    add_shape(slide3, x + Inches(0.1), impact_y, card_w - Inches(0.2), Inches(0.38),
              RGBColor(0xFD, 0xED, 0xED))
    impact_box = slide3.shapes.add_textbox(
        x + Inches(0.2), impact_y + Inches(0.04), card_w - Inches(0.4), Inches(0.3))
    tf = impact_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "IMPACT: "
    run1.font.size = Pt(7.5)
    run1.font.color.rgb = RED_ALERT
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = g["impact"]
    run2.font.size = Pt(7.5)
    run2.font.color.rgb = DARK_GRAY
    run2.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: Key Takeaways & Action Items
# ═══════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

# Header bar
add_shape(slide4, Inches(0), Inches(0), SLIDE_W, Inches(0.85), SAMSUNG_DARK)
add_text_box(slide4, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
             "Strategic Takeaways & Recommended Actions",
             font_size=22, color=WHITE, bold=True)

# Left column: Strategic Actions
add_text_box(slide4, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
             "PRIORITY ACTIONS FOR SAMSUNG",
             font_size=14, color=SAMSUNG_BLUE, bold=True)
add_shape(slide4, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.04), SAMSUNG_BLUE)

actions = [
    {
        "num": "1",
        "title": "Build Multi-Appliance Cooking Orchestration",
        "detail": "Develop Samsung's answer to Bosch Cook AI -- coordinate fridge, cooktop, and oven from a single AI system. Leverage SmartThings infrastructure for cross-appliance real-time control during cooking.",
        "color": RED_ALERT,
        "priority": "CRITICAL"
    },
    {
        "num": "2",
        "title": "Adopt Hybrid Edge-Cloud AI Architecture",
        "detail": "Keep camera/visual data on-device (match GE's privacy standard) while using Gemini cloud for generative features. Addresses 47% of users worried about data breaches.",
        "color": RED_ALERT,
        "priority": "HIGH"
    },
    {
        "num": "3",
        "title": "Enhance Samsung Food with Generative AI",
        "detail": "Add photo-to-recipe (restaurant dish → recipe), handwritten recipe digitization, and grocery delivery integration to match Flavorly AI capabilities.",
        "color": AMBER,
        "priority": "HIGH"
    },
    {
        "num": "4",
        "title": "Invest in Advanced Cooking Sensors",
        "detail": "Evaluate lambda-probe-equivalent sensors (Bosch) and smart cookware partnerships (Miele M Sense model). New sensor data improves AI cooking accuracy.",
        "color": AMBER,
        "priority": "MEDIUM"
    },
]

for i, action in enumerate(actions):
    y = Inches(1.7) + i * Inches(1.35)

    # Number circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = action["color"]
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = action["num"]
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Priority badge
    add_text_box(slide4, Inches(5), y + Inches(0.02), Inches(1.2), Inches(0.25),
                 action["priority"], font_size=8, color=action["color"], bold=True, alignment=PP_ALIGN.RIGHT)

    # Title
    add_text_box(slide4, Inches(1.05), y - Inches(0.02), Inches(4.5), Inches(0.3),
                 action["title"], font_size=11.5, color=DARK_GRAY, bold=True)

    # Detail
    add_text_box(slide4, Inches(1.05), y + Inches(0.3), Inches(5.2), Inches(0.85),
                 action["detail"], font_size=9, color=MED_GRAY)

# Right column: Category Leadership Summary
right_x = Inches(7)
add_text_box(slide4, right_x, Inches(1.1), Inches(6), Inches(0.4),
             "WHERE SAMSUNG LEADS vs TRAILS",
             font_size=14, color=SAMSUNG_BLUE, bold=True)
add_shape(slide4, right_x, Inches(1.5), Inches(5.8), Inches(0.04), SAMSUNG_BLUE)

# Samsung Leads section
add_shape(slide4, right_x, Inches(1.7), Inches(5.8), Inches(0.32), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide4, right_x + Inches(0.15), Inches(1.72), Inches(5.5), Inches(0.28),
             "SAMSUNG LEADS (5 of 11 categories)",
             font_size=10, color=GREEN, bold=True)

leads = [
    ("Fridge Food Recognition", "Google Gemini multimodal — most advanced AI partner"),
    ("Voice AI", "Only brand with Voice ID multi-user + built-in mic/speaker"),
    ("Energy Optimization", "Unique AI Hybrid Cooling (Peltier + compressor)"),
    ("Smart Home Ecosystem", "SmartThings unmatched cross-device breadth"),
    ("Security", "Knox Matrix hardware chip + blockchain — no competitor close"),
]

leads_box = slide4.shapes.add_textbox(right_x + Inches(0.15), Inches(2.1), Inches(5.5), Inches(2.2))
tf = leads_box.text_frame
tf.word_wrap = True
for i, (cat, desc) in enumerate(leads):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(3)
    run1 = p.add_run()
    run1.text = f"✓ {cat}: "
    run1.font.size = Pt(9)
    run1.font.color.rgb = GREEN
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(8.5)
    run2.font.color.rgb = MED_GRAY
    run2.font.name = "Calibri"

# Samsung Trails section
trails_y = Inches(4.15)
add_shape(slide4, right_x, trails_y, Inches(5.8), Inches(0.32), RGBColor(0xFD, 0xED, 0xED))
add_text_box(slide4, right_x + Inches(0.15), trails_y + Inches(0.02), Inches(5.5), Inches(0.28),
             "COMPETITORS LEAD (Samsung must close gaps)",
             font_size=10, color=RED_ALERT, bold=True)

trails = [
    ("Autonomous Cooking", "Bosch Cook AI", "multi-appliance orchestration"),
    ("Sensor Innovation", "Bosch (lambda probe) + Miele (M Sense cookware)", "unique sensors"),
    ("Recipe AI", "GE Flavorly AI", "most capable generative platform"),
    ("Maintenance AI", "Miele AI Diagnostics", "80%+ self-resolution rate"),
    ("Edge AI / Privacy", "GE/Haier", "images on-device, HomeGPT certified edge LLM"),
    ("Display Tech", "LG T-OLED", "36\" transparent OLED on fridge — industry first"),
]

trails_box = slide4.shapes.add_textbox(right_x + Inches(0.15), trails_y + Inches(0.38), Inches(5.5), Inches(2.5))
tf = trails_box.text_frame
tf.word_wrap = True
for i, (cat, leader, desc) in enumerate(trails):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(3)
    run1 = p.add_run()
    run1.text = f"✗ {cat}: "
    run1.font.size = Pt(9)
    run1.font.color.rgb = RED_ALERT
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = f"{leader} — {desc}"
    run2.font.size = Pt(8.5)
    run2.font.color.rgb = MED_GRAY
    run2.font.name = "Calibri"

# Bottom bar: forward-looking statement
bottom_y = Inches(6.85)
add_shape(slide4, Inches(0), bottom_y, SLIDE_W, Inches(0.55), RGBColor(0xE8, 0xED, 0xF5))

bottom_box = slide4.shapes.add_textbox(Inches(0.5), bottom_y + Inches(0.08), Inches(12.3), Inches(0.4))
tf = bottom_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "2026-2028 OUTLOOK: "
run1.font.size = Pt(10)
run1.font.color.rgb = SAMSUNG_DARK
run1.font.bold = True
run1.font.name = "Calibri"
run2 = p.add_run()
run2.text = ("Agentic AI cooking becomes standard for premium appliances  |  On-device LLMs proliferate (privacy)  |  "
             "Matter adoption erodes proprietary ecosystems  |  Smart cookware spawns new product category  |  "
             "Generative recipe AI becomes table stakes")
run2.font.size = Pt(9)
run2.font.color.rgb = MED_GRAY
run2.font.name = "Calibri"


# ── Save ──
output_path = "/home/chachapranu/ashu/git/product-research/claude/Samsung_AI_Kitchen_Competitive_Landscape_2026.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
