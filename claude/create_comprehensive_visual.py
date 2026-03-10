#!/usr/bin/env python3
"""
Samsung Kitchen: Comprehensive Product & Tech Competitive Landscape (Visual Variant)
4-slide presentation - less text, more visual impact, big numbers, clean layout
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)
SAMSUNG_DARK = RGBColor(0x0A, 0x12, 0x5E)
SAMSUNG_LIGHT = RGBColor(0xE8, 0xEB, 0xF7)
LG_RED = RGBColor(0xA5, 0x00, 0x34)
BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
GE_BLUE = RGBColor(0x00, 0x4B, 0x8D)
WHIRLPOOL_GOLD = RGBColor(0xC8, 0x96, 0x23)
MIELE_BROWN = RGBColor(0x6E, 0x1D, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xF3, 0x9C, 0x12)
RED_ALERT = RGBColor(0xE7, 0x4C, 0x3C)
TEAL = RGBColor(0x00, 0x89, 0x7B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: Title with big stats
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, SAMSUNG_DARK)

# Decorative circles
add_circle(slide1, Inches(-1), Inches(-1.5), Inches(4), RGBColor(0x10, 0x1D, 0x70))
add_circle(slide1, Inches(10.5), Inches(5), Inches(4), RGBColor(0x10, 0x1D, 0x70))
add_circle(slide1, Inches(8), Inches(-2), Inches(2.5), RGBColor(0x0D, 0x17, 0x68))

# Title
add_text_box(slide1, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Samsung Kitchen", 44, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1), Inches(1.75), Inches(11), Inches(0.6),
             "Product & Tech Competitive Landscape 2026", 28, RGBColor(0x8A, 0x9C, 0xE8), False, PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1), Inches(2.4), Inches(8), Inches(0.4),
             "Comprehensive Product Research  \u2022  Software, AI & Features  \u2022  Confidential", 13, RGBColor(0x6A, 0x7C, 0xC8), False)

# Big stat cards
stat_cards = [
    ("$311B", "Global Kitchen\nMarket (2025)", "$40-60B smart by 2030"),
    ("21.7%", "Samsung Smart\nKitchen Share", "#1 in fastest-growing segment"),
    ("4.55/5", "Composite KPI\nScore", "Leads all 6 competitors"),
    ("6", "Major Competitors\nAnalyzed", "LG, Bosch, GE, Whirlpool, Miele"),
]

card_w = Inches(2.6)
card_h = Inches(2.8)
card_gap = Inches(0.4)
start_x = Inches(1)
start_y = Inches(3.4)

for i, (big_num, label, sub) in enumerate(stat_cards):
    x = start_x + i * (card_w + card_gap)
    add_rounded_rect(slide1, x, start_y, card_w, card_h, RGBColor(0x12, 0x20, 0x80))
    add_text_box(slide1, x, start_y + Inches(0.4), card_w, Inches(0.8),
                 big_num, 36, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide1, x, start_y + Inches(1.3), card_w, Inches(0.6),
                 label, 14, RGBColor(0xB0, 0xBE, 0xE8), False, PP_ALIGN.CENTER)
    # Divider line
    add_shape(slide1, x + Inches(0.8), start_y + Inches(2.0), Inches(1.0), Inches(0.02), RGBColor(0x1A, 0x30, 0x90))
    add_text_box(slide1, x, start_y + Inches(2.15), card_w, Inches(0.4),
                 sub, 9, RGBColor(0x7A, 0x8C, 0xD0), False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Product Feature Showdown (Visual Cards)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# Thin top accent
add_shape(slide2, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SAMSUNG_BLUE)

add_text_box(slide2, Inches(0.5), Inches(0.2), Inches(8), Inches(0.45),
             "Product Feature Showdown", 26, SAMSUNG_DARK, True)
add_text_box(slide2, Inches(0.5), Inches(0.65), Inches(10), Inches(0.3),
             "Samsung vs 5 competitors across refrigeration, cooking, AI platforms, and hardware innovation", 11, MED_GRAY, False)

# Feature comparison cards - 3 columns x 4 rows
features = [
    # Row 1: Refrigeration
    ("\U0001F9CA", "Fridge Display", "Samsung", "32\" FHD Touch\n(Industry largest)", "LG: T-OLED 36\"\nGE: 7\" LCD", SAMSUNG_BLUE, GREEN),
    ("\U0001F4F7", "Fridge Camera AI", "Samsung", "Dual-lens + Gemini\n1000s of items", "LG: remote view\nGE: barcode scan", SAMSUNG_BLUE, GREEN),
    ("\u2744\uFE0F", "Cooling Tech", "Samsung", "AI Hybrid (Peltier)\nWorld-first dual", "LG: Inverter Linear\nMiele: MasterCool", SAMSUNG_BLUE, GREEN),
    # Row 2: Cooking
    ("\U0001F373", "Oven AI", "Bosch", "Cook AI: Agentic\nMulti-appliance", "Samsung: 80 dishes\nLG: 85+ dishes", BOSCH_RED, RED_ALERT),
    ("\U0001F321\uFE0F", "Cooking Sensors", "Bosch/Miele", "Lambda O2 probe\nM Sense cookware", "Samsung: camera\nLG: camera", BOSCH_RED, AMBER),
    ("\U0001F4E6", "Barcode Scanner", "GE", "Scan-to-List\n4M+ products", "Samsung: none\nLG: none", GE_BLUE, RED_ALERT),
    # Row 3: AI/Software
    ("\U0001F916", "LLM/GenAI", "Samsung", "Google Gemini\nMultimodal", "LG: EXAONE\nGE: Vertex AI", SAMSUNG_BLUE, GREEN),
    ("\U0001F3E0", "Ecosystem Breadth", "Samsung", "Phone+TV+Watch\n+Home (broadest)", "LG: TV+Home\nOthers: appliance only", SAMSUNG_BLUE, GREEN),
    ("\U0001F4F1", "Recipe Platform", "Samsung", "Samsung Food\n160K+ recipes", "Whirlpool: Yummly 27M\nGE: Flavorly AI", SAMSUNG_BLUE, AMBER),
    # Row 4: Trust/Innovation
    ("\U0001F3C6", "Brand Trust", "Bosch", "#1 most trusted\n7 consecutive years", "Samsung: ranked 4th\nLG: ranked 3rd", BOSCH_RED, RED_ALERT),
    ("\U0001F4C8", "Share Growth", "GE", "+2 ppt gain (2025)\nFastest growing", "Samsung: steady 18%\nLG: stable 17%", GE_BLUE, AMBER),
    ("\U0001F50C", "Matter Protocol", "Bosch", "First-mover\nShipped first device", "Samsung: 1.5 camera\nLG: Yes", BOSCH_RED, AMBER),
]

card_w = Inches(3.9)
card_h = Inches(1.3)
gap_x = Inches(0.35)
gap_y = Inches(0.15)
start_x = Inches(0.5)
start_y = Inches(1.1)

for i, (icon, title, leader, leader_detail, others, leader_color, samsung_status) in enumerate(features):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    # Card background
    add_rounded_rect(slide2, x, y, card_w, card_h, LIGHT_GRAY, RGBColor(0xE0, 0xE0, 0xE0))

    # Icon
    add_text_box(slide2, x + Inches(0.12), y + Inches(0.12), Inches(0.35), Inches(0.35), icon, 18, DARK_GRAY, False)

    # Title
    add_text_box(slide2, x + Inches(0.5), y + Inches(0.1), Inches(1.8), Inches(0.25), title, 11, DARK_GRAY, True)

    # Leader badge
    badge_color = leader_color
    add_rounded_rect(slide2, x + Inches(2.5), y + Inches(0.1), Inches(1.2), Inches(0.22), badge_color)
    add_text_box(slide2, x + Inches(2.5), y + Inches(0.1), Inches(1.2), Inches(0.22),
                 f"\u2605 {leader}", 8, WHITE, True, PP_ALIGN.CENTER)

    # Leader detail
    add_text_box(slide2, x + Inches(0.5), y + Inches(0.4), Inches(1.8), Inches(0.8), leader_detail, 8.5, DARK_GRAY, False)

    # Others
    add_text_box(slide2, x + Inches(2.4), y + Inches(0.4), Inches(1.4), Inches(0.8), others, 7.5, MED_GRAY, False)

    # Samsung status indicator (if Samsung doesn't lead)
    if leader != "Samsung":
        status_color = samsung_status
        status_text = "GAP" if status_color == RED_ALERT else "WATCH"
        add_rounded_rect(slide2, x + Inches(0.12), y + Inches(0.95), Inches(0.65), Inches(0.2), status_color)
        add_text_box(slide2, x + Inches(0.12), y + Inches(0.95), Inches(0.65), Inches(0.2),
                     f"Samsung: {status_text}", 6.5, WHITE, True, PP_ALIGN.CENTER)
    else:
        add_rounded_rect(slide2, x + Inches(0.12), y + Inches(0.95), Inches(0.65), Inches(0.2), GREEN)
        add_text_box(slide2, x + Inches(0.12), y + Inches(0.95), Inches(0.65), Inches(0.2),
                     "Samsung: LEAD", 6.5, WHITE, True, PP_ALIGN.CENTER)

# Summary strip
add_shape(slide2, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), SAMSUNG_DARK)
add_text_box(slide2, Inches(0.5), Inches(7.0), Inches(3), Inches(0.2),
             "SAMSUNG SCORECARD:", 10, WHITE, True)
summary_parts = [
    ("5 LEADS", GREEN),
    ("4 WATCH", AMBER),
    ("3 GAPS", RED_ALERT),
]
sx = Inches(3.0)
for text, color in summary_parts:
    add_rounded_rect(slide2, sx, Inches(7.0), Inches(1.1), Inches(0.25), color)
    add_text_box(slide2, sx, Inches(7.0), Inches(1.1), Inches(0.25), text, 9, WHITE, True, PP_ALIGN.CENTER)
    sx += Inches(1.3)

add_text_box(slide2, Inches(7), Inches(7.0), Inches(5.5), Inches(0.4),
             "Samsung leads in ecosystem + display + AI vision. Gaps in autonomous cooking, brand trust, and barcode innovation.", 8.5, RGBColor(0xB0, 0xBE, 0xE8), False)


# ============================================================
# SLIDE 3: AI Architecture & Tech Strategy
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_shape(slide3, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SAMSUNG_BLUE)

add_text_box(slide3, Inches(0.5), Inches(0.2), Inches(8), Inches(0.45),
             "AI Architecture & Tech Strategy", 26, SAMSUNG_DARK, True)
add_text_box(slide3, Inches(0.5), Inches(0.65), Inches(10), Inches(0.3),
             "Three competing philosophies shaping the autonomous kitchen era", 11, MED_GRAY, False)

# Three large architecture cards
arch_cards = [
    ("Samsung", "DISTRIBUTED EDGE", "Exynos NPU in every appliance", SAMSUNG_BLUE,
     "<50ms", "latency",
     "HIGH", "CapEx/unit",
     "Gemini", "cloud partner",
     [
         "NPU does local object detection",
         "Gemini handles complex reasoning",
         "Works offline for basic functions",
         "Hard to upgrade (chipset fixed)",
     ],
     "RISK: Cloud dependency for 15-20yr appliance lifespan"),
    ("LG", "CENTRAL HUB", "ThinQ ON brain, dumb endpoints", LG_RED,
     "100-300ms", "latency",
     "LOW", "CapEx/unit",
     "EXAONE", "on-device LLM",
     [
         "DQ-X AI Chipset in central hub",
         "1.2B-param SLM runs locally",
         "Easy to upgrade (swap hub)",
         "Single point of failure risk",
     ],
     "STRENGTH: On-device AI solves cloud dependency"),
    ("Bosch", "PHYSICS-FIRST", "Sensors over cameras", BOSCH_RED,
     "Real-time", "latency",
     "LOWEST", "CapEx/unit",
     "Internal", "AI engine",
     [
         "Lambda O2 sensor from automotive",
         "Humidity curves predict doneness",
         "No cloud needed for cooking AI",
         "Less \"visible\" AI to consumers",
     ],
     "STRENGTH: Best cooking results, cloud-independent"),
]

card_w = Inches(3.9)
card_h = Inches(4.6)
start_x = Inches(0.5)
gap = Inches(0.35)

for i, (brand, arch_name, tagline, color, stat1, label1, stat2, label2, stat3, label3, bullets, risk) in enumerate(arch_cards):
    x = start_x + i * (card_w + gap)
    y = Inches(1.1)

    # Card
    add_rounded_rect(slide3, x, y, card_w, card_h, LIGHT_GRAY, RGBColor(0xDD, 0xDD, 0xDD))

    # Brand header
    add_shape(slide3, x, y, card_w, Inches(0.7), color)
    add_text_box(slide3, x + Inches(0.2), y + Inches(0.08), card_w - Inches(0.4), Inches(0.3),
                 brand, 18, WHITE, True)
    add_text_box(slide3, x + Inches(0.2), y + Inches(0.38), card_w - Inches(0.4), Inches(0.25),
                 arch_name, 10, RGBColor(0xFF, 0xFF, 0xFF), True)

    # Tagline
    add_text_box(slide3, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.3),
                 tagline, 9, MED_GRAY, False, PP_ALIGN.CENTER)

    # 3 stat boxes
    stat_w = Inches(1.05)
    stat_h = Inches(0.65)
    stats = [(stat1, label1), (stat2, label2), (stat3, label3)]
    for j, (val, lbl) in enumerate(stats):
        sx = x + Inches(0.15) + j * (stat_w + Inches(0.1))
        add_rounded_rect(slide3, sx, y + Inches(1.15), stat_w, stat_h, WHITE)
        add_text_box(slide3, sx, y + Inches(1.2), stat_w, Inches(0.3), val, 14, color, True, PP_ALIGN.CENTER)
        add_text_box(slide3, sx, y + Inches(1.5), stat_w, Inches(0.2), lbl, 7, MED_GRAY, False, PP_ALIGN.CENTER)

    # Bullet points
    for j, bullet in enumerate(bullets):
        by = y + Inches(1.95) + j * Inches(0.35)
        add_text_box(slide3, x + Inches(0.25), by, card_w - Inches(0.5), Inches(0.3),
                     f"\u2022 {bullet}", 9, DARK_GRAY, False)

    # Risk/strength bar
    is_risk = "RISK" in risk
    bar_color = AMBER if is_risk else GREEN
    add_shape(slide3, x + Inches(0.15), y + Inches(3.55), card_w - Inches(0.3), Inches(0.03), bar_color)
    add_text_box(slide3, x + Inches(0.15), y + Inches(3.65), card_w - Inches(0.3), Inches(0.8),
                 risk, 8, MED_GRAY, False, PP_ALIGN.CENTER)

# KPI Comparison strip at bottom
add_shape(slide3, Inches(0), Inches(5.9), SLIDE_W, Inches(0.04), SAMSUNG_BLUE)
add_text_box(slide3, Inches(0.5), Inches(6.05), Inches(3), Inches(0.3),
             "COMPOSITE KPI SCORES", 12, SAMSUNG_DARK, True)

# Visual KPI bars
kpi_brands = [
    ("Samsung", 4.55, SAMSUNG_BLUE),
    ("LG", 3.96, LG_RED),
    ("Bosch", 3.48, BOSCH_RED),
    ("GE/Haier", 3.05, GE_BLUE),
    ("Whirlpool", 2.88, WHIRLPOOL_GOLD),
    ("Miele", 2.57, MIELE_BROWN),
]

bar_start_x = Inches(0.5)
bar_max_w = Inches(8.0)
bar_h = Inches(0.16)
bar_gap = Inches(0.05)

for i, (brand, score, color) in enumerate(kpi_brands):
    by = Inches(6.35) + i * (bar_h + bar_gap)
    # Brand label
    add_text_box(slide3, bar_start_x, by, Inches(1.0), bar_h, brand, 8, DARK_GRAY, True, PP_ALIGN.RIGHT)
    # Bar background
    add_shape(slide3, Inches(1.6), by, bar_max_w, bar_h, RGBColor(0xE8, 0xE8, 0xE8))
    # Filled bar
    fill_w = Emu(int(bar_max_w * score / 5.0))
    add_shape(slide3, Inches(1.6), by, fill_w, bar_h, color)
    # Score label
    add_text_box(slide3, Inches(1.6) + fill_w + Inches(0.1), by, Inches(0.5), bar_h,
                 str(score), 8, color, True)

# Right side: dimension breakdown
add_text_box(slide3, Inches(10.2), Inches(6.05), Inches(3), Inches(0.3),
             "SAMSUNG BY DIMENSION", 10, SAMSUNG_DARK, True)

dims = [
    ("Technology", "5.0", GREEN),
    ("Innovation", "5.0", GREEN),
    ("Product", "4.75", GREEN),
    ("Business", "4.75", GREEN),
    ("Market", "3.7", AMBER),
]

for i, (dim, score, color) in enumerate(dims):
    dy = Inches(6.35) + i * Inches(0.2)
    add_text_box(slide3, Inches(10.2), dy, Inches(1.2), Inches(0.2), dim, 8, DARK_GRAY, False)
    add_rounded_rect(slide3, Inches(11.4), dy, Inches(0.55), Inches(0.18), color)
    add_text_box(slide3, Inches(11.4), dy, Inches(0.55), Inches(0.18), score, 8, WHITE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: Strategic Gaps & Action Plan
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_shape(slide4, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SAMSUNG_BLUE)

add_text_box(slide4, Inches(0.5), Inches(0.2), Inches(8), Inches(0.45),
             "Strategic Gaps & Action Plan", 26, SAMSUNG_DARK, True)
add_text_box(slide4, Inches(0.5), Inches(0.65), Inches(10), Inches(0.3),
             "7 priorities to maintain Samsung's competitive leadership in the evolving kitchen landscape", 11, MED_GRAY, False)

# Top row: 3 HIGH priority gap cards
high_gaps = [
    ("1", "Close Trust &\nReliability Gap", "vs Bosch (7yr #1 trust)",
     "#4 in trust\n87 problems/100 units\nApp rating: 3.8 vs 4.5 target",
     "Reliability Guarantee program\nPublish transparent data\nInvest in connected QA",
     BOSCH_RED, RED_ALERT),
    ("2", "Build Autonomous\nCooking AI", "vs Bosch Cook AI (agentic)",
     "Samsung is vision-focused\nBosch orchestrates multi-appliance\nLambda sensor outperforms camera",
     "Develop agentic cooking AI\nFridge\u2192Recipe\u2192Oven automation\nLeverage Gemini multimodal",
     BOSCH_RED, RED_ALERT),
    ("3", "Defend Smart Share\nvs GE Momentum", "vs GE (+2 ppt, 9x awards)",
     "GE: fastest share growth\nBarcode scanner unique HW\nInstacart grocery commerce",
     "Match barcode innovation\nDeepen grocery integration\nLeverage phone+TV ecosystem",
     GE_BLUE, RED_ALERT),
]

hcard_w = Inches(3.9)
hcard_h = Inches(2.8)
hcard_gap = Inches(0.35)
hcard_y = Inches(1.1)

for i, (num, title, vs_text, gap_text, action_text, brand_color, sev_color) in enumerate(high_gaps):
    x = Inches(0.5) + i * (hcard_w + hcard_gap)

    add_rounded_rect(slide4, x, hcard_y, hcard_w, hcard_h, LIGHT_GRAY, RGBColor(0xDD, 0xDD, 0xDD))

    # Number + severity header
    add_shape(slide4, x, hcard_y, hcard_w, Inches(0.4), sev_color)
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), hcard_y + Inches(0.05), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    add_text_box(slide4, x + Inches(0.1), hcard_y + Inches(0.07), Inches(0.3), Inches(0.25),
                 num, 12, sev_color, True, PP_ALIGN.CENTER)
    add_text_box(slide4, x + Inches(0.5), hcard_y + Inches(0.08), Inches(1.5), Inches(0.25),
                 "HIGH PRIORITY", 9, WHITE, True)

    # Title
    add_text_box(slide4, x + Inches(0.15), hcard_y + Inches(0.5), hcard_w - Inches(0.3), Inches(0.5),
                 title, 13, DARK_GRAY, True)

    # vs badge
    add_rounded_rect(slide4, x + Inches(0.15), hcard_y + Inches(1.05), hcard_w - Inches(0.3), Inches(0.22), brand_color)
    add_text_box(slide4, x + Inches(0.15), hcard_y + Inches(1.05), hcard_w - Inches(0.3), Inches(0.22),
                 vs_text, 7.5, WHITE, True, PP_ALIGN.CENTER)

    # Gap | Action split
    half_w = (hcard_w - Inches(0.4)) / 2

    add_text_box(slide4, x + Inches(0.15), hcard_y + Inches(1.35), half_w, Inches(0.18),
                 "THE GAP", 7, RED_ALERT, True)
    add_text_box(slide4, x + Inches(0.15), hcard_y + Inches(1.55), half_w, Inches(1.1),
                 gap_text, 8, DARK_GRAY, False)

    add_text_box(slide4, x + Inches(0.15) + half_w + Inches(0.1), hcard_y + Inches(1.35), half_w, Inches(0.18),
                 "ACTION", 7, GREEN, True)
    add_text_box(slide4, x + Inches(0.15) + half_w + Inches(0.1), hcard_y + Inches(1.55), half_w, Inches(1.1),
                 action_text, 8, DARK_GRAY, False)

# Bottom row: 4 MEDIUM priority cards
med_gaps = [
    ("4", "Cloud Dependency", "MED-HIGH",
     "Gemini 15-20yr risk\nLG/Bosch run local",
     "Hybrid AI: core on-device\nGraceful cloud degradation", AMBER),
    ("5", "KitchenAid Relaunch", "MED-HIGH",
     "10yr biggest relaunch\nCamera AI + 27M Yummly",
     "Differentiate on ecosystem\nDeepen SmartThings + Gemini", AMBER),
    ("6", "Dacor Luxury Gap", "MEDIUM",
     "Below Thermador/SubZero\nMiele 20yr standard",
     "Dacor x Samsung AI suite\nAI luxury differentiation", AMBER),
    ("7", "Chinese Brand Entry", "MEDIUM",
     "Midea +17.7% growth\nMidea-Hisense AI pact",
     "Bespoke entry-tier expand\nEcosystem moat vs price", AMBER),
]

mcard_w = Inches(2.9)
mcard_h = Inches(2.55)
mcard_gap = Inches(0.3)
mcard_y = Inches(4.15)

for i, (num, title, severity, gap_text, action_text, sev_color) in enumerate(med_gaps):
    x = Inches(0.5) + i * (mcard_w + mcard_gap)

    add_rounded_rect(slide4, x, mcard_y, mcard_w, mcard_h, LIGHT_GRAY, RGBColor(0xDD, 0xDD, 0xDD))

    # Header
    add_shape(slide4, x, mcard_y, mcard_w, Inches(0.35), sev_color)
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.08), mcard_y + Inches(0.04), Inches(0.27), Inches(0.27))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    add_text_box(slide4, x + Inches(0.08), mcard_y + Inches(0.05), Inches(0.27), Inches(0.22),
                 num, 10, sev_color, True, PP_ALIGN.CENTER)
    add_text_box(slide4, x + Inches(0.4), mcard_y + Inches(0.07), Inches(1.8), Inches(0.22),
                 severity, 8, WHITE, True)

    # Title
    add_text_box(slide4, x + Inches(0.12), mcard_y + Inches(0.45), mcard_w - Inches(0.24), Inches(0.3),
                 title, 12, DARK_GRAY, True)

    # Gap
    add_text_box(slide4, x + Inches(0.12), mcard_y + Inches(0.8), mcard_w - Inches(0.24), Inches(0.15),
                 "GAP", 7, RED_ALERT, True)
    add_text_box(slide4, x + Inches(0.12), mcard_y + Inches(0.95), mcard_w - Inches(0.24), Inches(0.6),
                 gap_text, 8, DARK_GRAY, False)

    # Action
    add_text_box(slide4, x + Inches(0.12), mcard_y + Inches(1.55), mcard_w - Inches(0.24), Inches(0.15),
                 "ACTION", 7, GREEN, True)
    add_text_box(slide4, x + Inches(0.12), mcard_y + Inches(1.7), mcard_w - Inches(0.24), Inches(0.7),
                 action_text, 8, DARK_GRAY, False)

# Bottom outlook bar
add_shape(slide4, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), SAMSUNG_DARK)
add_text_box(slide4, Inches(0.5), Inches(7.0), Inches(2.5), Inches(0.25),
             "2026-2030 OUTLOOK", 11, WHITE, True)
timeline_items = [
    ("2026", "Smart kitchen\ntipping point"),
    ("2027", "Autonomous cooking\ngoes commercial"),
    ("2028", "25% smart\npenetration"),
    ("2029", "Fully agentic\nkitchens (premium)"),
    ("2030", "$40-60B smart\nkitchen market"),
]

for i, (year, desc) in enumerate(timeline_items):
    tx = Inches(3.2) + i * Inches(2.0)
    # Year badge
    add_rounded_rect(slide4, tx, Inches(6.98), Inches(0.55), Inches(0.2), SAMSUNG_BLUE)
    add_text_box(slide4, tx, Inches(6.98), Inches(0.55), Inches(0.2), year, 8, WHITE, True, PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide4, tx + Inches(0.6), Inches(6.98), Inches(1.3), Inches(0.45), desc, 7.5, RGBColor(0xB0, 0xBE, 0xE8), False)


# ============================================================
# SLIDE 5: Appendix - Data Sources (Visual Style)
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

# Thin top accent
add_shape(slide5, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SAMSUNG_BLUE)

add_text_box(slide5, Inches(0.5), Inches(0.2), Inches(8), Inches(0.45),
             "Appendix: Data Sources", 26, SAMSUNG_DARK, True)
add_text_box(slide5, Inches(0.5), Inches(0.65), Inches(10), Inches(0.3),
             "100+ primary sources from market research firms, company newsrooms, and CES/KBIS 2026 coverage", 11, MED_GRAY, False)

# Source category cards - 3 columns x 2 rows
source_cards = [
    ("\U0001F4CA", "Market Data & Sizing", SAMSUNG_BLUE, [
        "$311B market, CAGR 4.2-4.8%\n\u2192 Mordor Intelligence, IMARC, Coherent MI",
        "Smart kitchen $23.7B \u2192 $40-60B\n\u2192 Research and Markets, Statista, Polaris",
        "US smart $6.9B \u2192 $17.64B by 2033\n\u2192 GlobeNewsWire US Forecast 2025-2033",
        "13% \u2192 30.8% smart penetration\n\u2192 Statista Digital Market Insights",
    ]),
    ("\U0001F3C6", "Market Share & Trust", GREEN, [
        "Samsung 18% / 21.7% smart share\n\u2192 OpenBrand Q1-Q2 2025 Reports",
        "Brand rankings (Whirlpool/LG/GE/Bosch)\n\u2192 OpenBrand, Fred's Appliance Academy",
        "Bosch #1 trust (7 years)\n\u2192 Lifestory Research 2025 Study",
        "87/100 connected problems, -11pt NPS\n\u2192 J.D. Power 2025 Appliance Study",
    ]),
    ("\U0001F4F1", "Samsung Product & AI", SAMSUNG_BLUE, [
        "Bespoke AI, Gemini, Hybrid Cooling\n\u2192 Samsung Newsroom, CES 2026 press",
        "SmartThings 425M+, Matter 1.5\n\u2192 SmartThings Blog, Samsung SEC filings",
        "Samsung Food 160K+ recipes\n\u2192 Samsung Food platform data",
        "Knox Matrix, 9,304 US patents\n\u2192 GreyB Analytics, Samsung filings",
    ]),
    ("\U0001F50D", "Competitor Data", MED_GRAY, [
        "LG: Gourmet AI, EXAONE, ThinQ ON\n\u2192 LG Newsroom CES 2026 releases",
        "Bosch: Cook AI, lambda probe\n\u2192 BSH CES 2026 (BusinessWire), Reviewed",
        "GE: Barcode scanner, Flavorly AI\n\u2192 GE Press Room, Google Cloud release",
        "Whirlpool/Miele: KA Camera, M Sense\n\u2192 KBIS 2026 (PRNewswire), Miele SBID",
    ]),
    ("\u2699\uFE0F", "Technology & Architecture", RGBColor(0x8E, 0x44, 0xAD), [
        "Edge vs Hub vs Sensor architectures\n\u2192 Gemini Market Intelligence Report",
        "Matter 1.2-1.5, 750+ certified products\n\u2192 Connectivity Standards Alliance (CSA)",
        "AI patent growth 800%+ since 2017\n\u2192 GreyB Patent Analytics",
        "47% data breach concern\n\u2192 American Home Shield Smart Survey",
    ]),
    ("\U0001F3AA", "CES/KBIS 2026 Events", RGBColor(0xE6, 0x7E, 0x22), [
        "CES 2026 Innovation Awards\n\u2192 CES official, Reviewed, Tom's Guide",
        "KBIS 2026 Best of Show\n\u2192 KBIS official, PRNewswire",
        "GE 9x Smart Appliance of Year\n\u2192 GE Appliances Press Room",
        "239 smart kitchen startups\n\u2192 Tracxn Smart Kitchen Report",
    ]),
]

card_w = Inches(3.9)
card_h = Inches(2.65)
gap_x = Inches(0.35)
gap_y = Inches(0.2)
start_x = Inches(0.5)
start_y = Inches(1.1)

for i, (icon, title, color, items) in enumerate(source_cards):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    # Card
    add_rounded_rect(slide5, x, y, card_w, card_h, LIGHT_GRAY, RGBColor(0xE0, 0xE0, 0xE0))

    # Header bar
    add_shape(slide5, x, y, card_w, Inches(0.4), color)
    add_text_box(slide5, x + Inches(0.12), y + Inches(0.07), Inches(0.35), Inches(0.3), icon, 16, WHITE, False)
    add_text_box(slide5, x + Inches(0.45), y + Inches(0.08), card_w - Inches(0.6), Inches(0.28), title, 11, WHITE, True)

    # Source items
    for j, item in enumerate(items):
        iy = y + Inches(0.5) + j * Inches(0.52)
        lines = item.split("\n")
        # Data point
        add_text_box(slide5, x + Inches(0.15), iy, card_w - Inches(0.3), Inches(0.2),
                     lines[0], 8, DARK_GRAY, True)
        # Source
        if len(lines) > 1:
            add_text_box(slide5, x + Inches(0.15), iy + Inches(0.18), card_w - Inches(0.3), Inches(0.2),
                         lines[1], 7, MED_GRAY, False)

# Footer
add_shape(slide5, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), SAMSUNG_DARK)
add_text_box(slide5, Inches(0.5), Inches(7.03), Inches(12), Inches(0.35),
             "Report compiled March 2026  \u2022  Publicly available data  \u2022  KPI scoring 1-5 scale  \u2022  "
             "Full source list (60+ URLs) in companion research document",
             8.5, RGBColor(0xB0, 0xBE, 0xE8), False)


# Save
output_path = "/home/chachapranu/ashu/git/product-research/claude/Samsung_Kitchen_Comprehensive_Visual_2026.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
