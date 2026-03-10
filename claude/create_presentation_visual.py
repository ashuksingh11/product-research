#!/usr/bin/env python3
"""
Samsung vs Competition: AI Kitchen Features - VISUAL VARIANT
4-Slide Executive Presentation — Minimal text, bold visuals, Samsung Blue branding
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
SAMSUNG_BLUE = RGBColor(0x14, 0x28, 0xA0)
SAMSUNG_LIGHT = RGBColor(0x3A, 0x5B, 0xDB)
SAMSUNG_DARK = RGBColor(0x0A, 0x12, 0x5E)
SAMSUNG_ACCENT = RGBColor(0x6B, 0x8C, 0xF2)
LG_RED = RGBColor(0xA5, 0x00, 0x34)
BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
GE_BLUE = RGBColor(0x00, 0x6D, 0xAA)
WHIRLPOOL_GOLD = RGBColor(0xC8, 0x96, 0x23)
MIELE_BROWN = RGBColor(0x8B, 0x2E, 0x3B)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFC)
BLACK = RGBColor(0x12, 0x12, 0x12)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
MED_GRAY = RGBColor(0x78, 0x78, 0x78)
LIGHT_GRAY = RGBColor(0xEE, 0xEF, 0xF3)
SOFT_BLUE = RGBColor(0xE3, 0xEA, 0xFA)
GREEN = RGBColor(0x1B, 0x8A, 0x4E)
SOFT_GREEN = RGBColor(0xE2, 0xF5, 0xEB)
RED_ALERT = RGBColor(0xD4, 0x2B, 0x2B)
SOFT_RED = RGBColor(0xFD, 0xE8, 0xE8)
AMBER = RGBColor(0xE6, 0x8A, 0x00)
SOFT_AMBER = RGBColor(0xFE, 0xF3, 0xE2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)
SH = Inches(7.5)


def shape(slide, left, top, w, h, fill=None, line=None, lw=None, rounding=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounding else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, left, top, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.fill.solid()
        s.line.color.rgb = line
        if lw:
            s.line.width = Pt(lw)
    return s


def circle(slide, left, top, size, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, left, top, w, h, text, sz=12, clr=BLACK, bold=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb


def rich_text(slide, left, top, w, h, runs_list, align=PP_ALIGN.LEFT):
    """runs_list: [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    for text, sz, clr, bold in runs_list:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = clr
        r.font.bold = bold
        r.font.name = "Calibri"
    return tb


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title — Bold Visual
# ═══════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])

# Full dark blue background
shape(s1, Inches(0), Inches(0), SW, SH, SAMSUNG_DARK)

# Decorative gradient-like accent shapes
shape(s1, Inches(0), Inches(0), Inches(0.12), SH, SAMSUNG_BLUE)
circle(s1, Inches(10.2), Inches(-1.5), Inches(5), RGBColor(0x16, 0x2E, 0xB0))
circle(s1, Inches(11), Inches(4.5), Inches(4.5), RGBColor(0x12, 0x22, 0x8A))

# Small Samsung accent bar
shape(s1, Inches(1.2), Inches(1.8), Inches(0.8), Inches(0.06), SAMSUNG_ACCENT)

# Title
txt(s1, Inches(1.2), Inches(2.0), Inches(8), Inches(1.2),
    "AI Kitchen", sz=48, clr=WHITE, bold=True)
txt(s1, Inches(1.2), Inches(2.85), Inches(8), Inches(1),
    "Competitive Landscape", sz=40, clr=SAMSUNG_ACCENT, bold=False)

# Subtitle
txt(s1, Inches(1.2), Inches(4.1), Inches(7), Inches(0.5),
    "Samsung vs LG  ·  Bosch  ·  GE/Haier  ·  Whirlpool  ·  Miele",
    sz=16, clr=RGBColor(0x9A, 0xAA, 0xDD))

# Divider line
shape(s1, Inches(1.2), Inches(4.75), Inches(4.5), Inches(0.02), RGBColor(0x3A, 0x50, 0xAA))

# Key stat boxes — bottom left
stats = [
    ("$17.6B", "US Smart Kitchen\nMarket by 2033"),
    ("21.7%", "Samsung Market\nShare"),
    ("6", "Major Brands\nCompeting in AI"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(1.2) + i * Inches(2.6)
    y = Inches(5.2)

    txt(s1, x, y, Inches(2.2), Inches(0.5), num, sz=28, clr=WHITE, bold=True)
    txt(s1, x, y + Inches(0.55), Inches(2.2), Inches(0.6), label, sz=10, clr=RGBColor(0x9A, 0xAA, 0xDD))

# Date
txt(s1, Inches(1.2), Inches(6.8), Inches(4), Inches(0.3),
    "March 2026  |  Product Research Deep Dive", sz=9, clr=RGBColor(0x6A, 0x7A, 0xBB))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Visual Comparison — Bubble / Radar Style
# ═══════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
shape(s2, Inches(0), Inches(0), SW, SH, OFF_WHITE)

# Header
shape(s2, Inches(0), Inches(0), SW, Inches(0.9), SAMSUNG_BLUE)
txt(s2, Inches(0.6), Inches(0.18), Inches(8), Inches(0.55),
    "Who Leads in Each AI Category?", sz=24, clr=WHITE, bold=True)

# Grid of category cards — 4 cols × 3 rows
categories = [
    # (category, leader, leader_color, samsung_position, icon)
    ("Fridge Vision AI", "SAMSUNG", SAMSUNG_BLUE, "1st", "👁"),
    ("Oven Cooking AI", "LG", LG_RED, "3rd", "🍳"),
    ("Autonomous Cooking", "BOSCH", BOSCH_RED, "4th", "🤖"),
    ("Voice & Conv. AI", "SAMSUNG", SAMSUNG_BLUE, "1st", "🎤"),
    ("Energy Optimization", "SAMSUNG", SAMSUNG_BLUE, "1st", "⚡"),
    ("Recipe AI Platform", "GE/HAIER", GE_BLUE, "2nd", "📖"),
    ("Predictive Maint.", "MIELE", MIELE_BROWN, "2nd", "🔧"),
    ("Sensor Innovation", "BOSCH + MIELE", BOSCH_RED, "3rd", "📡"),
    ("Ecosystem Breadth", "SAMSUNG", SAMSUNG_BLUE, "1st", "🏠"),
    ("Edge AI / Privacy", "GE/HAIER", GE_BLUE, "4th", "🔒"),
    ("Security Framework", "SAMSUNG", SAMSUNG_BLUE, "1st", "🛡"),
    ("Display Technology", "LG", LG_RED, "2nd", "📺"),
]

card_w = Inches(3.0)
card_h = Inches(1.7)
start_x = Inches(0.35)
start_y = Inches(1.15)
gap_x = Inches(0.18)
gap_y = Inches(0.18)

for i, (cat, leader, ldr_clr, sam_pos, icon) in enumerate(categories):
    col = i % 4
    row = i // 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    is_samsung_lead = leader == "SAMSUNG"
    card_bg = SOFT_BLUE if is_samsung_lead else WHITE
    border_clr = SAMSUNG_BLUE if is_samsung_lead else RGBColor(0xDD, 0xDD, 0xDD)

    # Card
    card = shape(s2, x, y, card_w, card_h, card_bg, line=border_clr, lw=1.5, rounding=True)

    # Icon
    txt(s2, x + Inches(0.15), y + Inches(0.1), Inches(0.4), Inches(0.4),
        icon, sz=20, align=PP_ALIGN.CENTER)

    # Category name
    txt(s2, x + Inches(0.55), y + Inches(0.12), Inches(2.3), Inches(0.35),
        cat, sz=11, clr=DARK_GRAY, bold=True)

    # Leader badge
    badge_bg = ldr_clr
    badge = shape(s2, x + Inches(0.15), y + Inches(0.6), Inches(2.7), Inches(0.35), badge_bg, rounding=True)
    badge_tb = s2.shapes.add_textbox(x + Inches(0.15), y + Inches(0.62), Inches(2.7), Inches(0.35))
    badge_tb.text_frame.word_wrap = True
    bp = badge_tb.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    r = bp.add_run()
    r.text = f"LEADER: {leader}"
    r.font.size = Pt(9)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Calibri"

    # Samsung position
    if not is_samsung_lead:
        pos_clr = AMBER if sam_pos in ("2nd",) else RED_ALERT
        pos_bg = SOFT_AMBER if sam_pos in ("2nd",) else SOFT_RED
        shape(s2, x + Inches(0.15), y + Inches(1.1), Inches(2.7), Inches(0.4), pos_bg, rounding=True)
        rich_text(s2, x + Inches(0.15), y + Inches(1.15), Inches(2.7), Inches(0.3),
                  [("Samsung: ", 9, MED_GRAY, False), (sam_pos, 11, pos_clr, True)],
                  align=PP_ALIGN.CENTER)
    else:
        # Checkmark for Samsung leads
        shape(s2, x + Inches(0.15), y + Inches(1.1), Inches(2.7), Inches(0.4), SOFT_GREEN, rounding=True)
        rich_text(s2, x + Inches(0.15), y + Inches(1.15), Inches(2.7), Inches(0.3),
                  [("✓  Samsung Leads", 10, GREEN, True)],
                  align=PP_ALIGN.CENTER)

# Bottom summary bar
shape(s2, Inches(0), Inches(6.85), SW, Inches(0.55), SAMSUNG_DARK)
rich_text(s2, Inches(0.6), Inches(6.93), Inches(12), Inches(0.35),
          [("Samsung leads 5 of 12 categories", 12, WHITE, True),
           ("   —   Trails in Autonomous Cooking (Bosch), Sensors (Bosch/Miele), Recipe AI (GE), Privacy (GE/Haier), Display (LG)", 10, RGBColor(0xAA, 0xBB, 0xEE), False)])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: Top 4 Competitive Gaps — Big Visual Cards
# ═══════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
shape(s3, Inches(0), Inches(0), SW, SH, OFF_WHITE)

# Header
shape(s3, Inches(0), Inches(0), SW, Inches(0.9), RGBColor(0x7A, 0x0C, 0x0C))
txt(s3, Inches(0.6), Inches(0.18), Inches(8), Inches(0.55),
    "Top Competitive Gaps Samsung Must Close", sz=24, clr=WHITE, bold=True)

# 4 big gap cards — 2x2 layout
gap_cards = [
    {
        "num": "01",
        "title": "No Multi-Appliance\nCooking AI",
        "competitor": "BOSCH COOK AI",
        "comp_color": BOSCH_RED,
        "what": "Coordinates cooktop + oven + probe simultaneously",
        "vs_samsung": "Samsung AI Pro Cooking is single-appliance only",
        "risk": "HIGH",
    },
    {
        "num": "02",
        "title": "Cloud-Dependent\nVision AI",
        "competitor": "GE EDGE AI",
        "comp_color": GE_BLUE,
        "what": "\"Images stay on device, not in the cloud\"",
        "vs_samsung": "Samsung sends food photos to Google Gemini cloud",
        "risk": "HIGH",
    },
    {
        "num": "03",
        "title": "No Smart\nCookware",
        "competitor": "MIELE M SENSE",
        "comp_color": MIELE_BROWN,
        "what": "3 temp sensors IN cookware + auto hob control",
        "vs_samsung": "Samsung has zero smart cookware strategy",
        "risk": "MEDIUM",
    },
    {
        "num": "04",
        "title": "Weaker Recipe\nAI Platform",
        "competitor": "GE FLAVORLY AI",
        "comp_color": GE_BLUE,
        "what": "Photo-to-recipe + handwritten digitization + Instacart",
        "vs_samsung": "Samsung Food lacks generative features",
        "risk": "MEDIUM",
    },
]

big_card_w = Inches(6.25)
big_card_h = Inches(2.8)
cx = Inches(0.25)
cy = Inches(1.15)
gx = Inches(0.2)
gy = Inches(0.2)

for i, g in enumerate(gap_cards):
    col = i % 2
    row = i // 2
    x = cx + col * (big_card_w + gx)
    y = cy + row * (big_card_h + gy)

    # Card background
    shape(s3, x, y, big_card_w, big_card_h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD), lw=1, rounding=True)

    # Risk indicator top-right
    risk_clr = RED_ALERT if g["risk"] == "HIGH" else AMBER
    risk_bg = SOFT_RED if g["risk"] == "HIGH" else SOFT_AMBER
    risk_w = Inches(1.0)
    shape(s3, x + big_card_w - risk_w - Inches(0.15), y + Inches(0.15), risk_w, Inches(0.3), risk_bg, rounding=True)
    txt(s3, x + big_card_w - risk_w - Inches(0.15), y + Inches(0.17), risk_w, Inches(0.28),
        f"● {g['risk']}", sz=9, clr=risk_clr, bold=True, align=PP_ALIGN.CENTER)

    # Large number
    txt(s3, x + Inches(0.25), y + Inches(0.15), Inches(0.8), Inches(0.7),
        g["num"], sz=36, clr=RGBColor(0xDD, 0xDD, 0xDD), bold=True)

    # Gap title — big and bold
    txt(s3, x + Inches(1.0), y + Inches(0.15), Inches(4), Inches(0.9),
        g["title"], sz=18, clr=DARK_GRAY, bold=True)

    # Competitor badge
    comp_badge_w = Inches(2.8)
    shape(s3, x + Inches(0.25), y + Inches(1.2), comp_badge_w, Inches(0.35), g["comp_color"], rounding=True)
    txt(s3, x + Inches(0.25), y + Inches(1.22), comp_badge_w, Inches(0.33),
        g["competitor"], sz=10, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # What they do
    rich_text(s3, x + Inches(0.25), y + Inches(1.7), Inches(5.7), Inches(0.4),
              [("What they do: ", 10, MED_GRAY, True), (g["what"], 10, DARK_GRAY, False)])

    # vs Samsung
    rich_text(s3, x + Inches(0.25), y + Inches(2.15), Inches(5.7), Inches(0.4),
              [("Samsung gap: ", 10, RED_ALERT, True), (g["vs_samsung"], 10, DARK_GRAY, False)])

# Also-watch strip at bottom
shape(s3, Inches(0), Inches(7.0), SW, Inches(0.45), SAMSUNG_DARK)

also_watch = [
    ("LG EXAONE", "On-device LLM in fridge", LG_RED),
    ("Bosch Matter-First", "Open ecosystem eroding SmartThings lock-in", BOSCH_RED),
    ("KitchenAid Relaunch", "Camera AI in premium ovens", WHIRLPOOL_GOLD),
    ("Midea", "Chinese brand disruption at lower prices", RGBColor(0x00, 0xB4, 0xD8)),
]

txt(s3, Inches(0.5), Inches(7.05), Inches(1.5), Inches(0.35),
    "ALSO WATCH:", sz=9, clr=RGBColor(0xAA, 0xBB, 0xEE), bold=True)

for i, (name, desc, clr) in enumerate(also_watch):
    ax = Inches(2.2) + i * Inches(2.8)
    # Small color dot
    c = circle(s3, ax, Inches(7.12), Inches(0.15), clr)
    rich_text(s3, ax + Inches(0.22), Inches(7.05), Inches(2.5), Inches(0.35),
              [(name, 9, WHITE, True), (f"  {desc}", 8, RGBColor(0x99, 0xAA, 0xCC), False)])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: Action Plan — Clean & Bold
# ═══════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
shape(s4, Inches(0), Inches(0), SW, SH, OFF_WHITE)

# Header
shape(s4, Inches(0), Inches(0), SW, Inches(0.9), SAMSUNG_BLUE)
txt(s4, Inches(0.6), Inches(0.18), Inches(8), Inches(0.55),
    "Samsung: 4 Strategic Priorities", sz=24, clr=WHITE, bold=True)

# 4 action columns
actions = [
    {
        "num": "1",
        "title": "Build Agentic\nCooking AI",
        "subtitle": "Multi-appliance orchestration",
        "points": [
            "Coordinate fridge + cooktop + oven from single AI",
            "Match Bosch Cook AI's autonomous cooking",
            "Leverage SmartThings for cross-device control",
        ],
        "color": SAMSUNG_BLUE,
        "priority": "CRITICAL",
        "pri_color": RED_ALERT,
    },
    {
        "num": "2",
        "title": "Keep Visual\nData On-Device",
        "subtitle": "Hybrid edge-cloud architecture",
        "points": [
            "Camera images processed locally (match GE)",
            "Use Gemini cloud only for generative features",
            "Address 47% user privacy concerns",
        ],
        "color": SAMSUNG_LIGHT,
        "priority": "HIGH",
        "pri_color": RED_ALERT,
    },
    {
        "num": "3",
        "title": "Upgrade\nSamsung Food",
        "subtitle": "Generative recipe capabilities",
        "points": [
            "Add photo-to-recipe (dish → recipe)",
            "Handwritten recipe digitization",
            "Integrate grocery delivery (Instacart-style)",
        ],
        "color": RGBColor(0x2E, 0x7D, 0x32),
        "priority": "HIGH",
        "pri_color": AMBER,
    },
    {
        "num": "4",
        "title": "Invest in\nSensor Tech",
        "subtitle": "Cookware & advanced sensors",
        "points": [
            "Evaluate smart cookware partnerships",
            "Add oven chemical sensors (lambda probe equiv.)",
            "Better sensor data → better AI cooking",
        ],
        "color": RGBColor(0x6A, 0x1B, 0x9A),
        "priority": "MEDIUM",
        "pri_color": AMBER,
    },
]

col_w = Inches(3.05)
col_h = Inches(4.8)
col_start_x = Inches(0.3)
col_start_y = Inches(1.15)
col_gap = Inches(0.15)

for i, a in enumerate(actions):
    x = col_start_x + i * (col_w + col_gap)
    y = col_start_y

    # Card
    shape(s4, x, y, col_w, col_h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD), lw=1, rounding=True)

    # Color top stripe
    shape(s4, x, y, col_w, Inches(0.08), a["color"], rounding=False)

    # Large number
    big_circle = circle(s4, x + Inches(1.1), y + Inches(0.3), Inches(0.7), a["color"])
    num_tb = s4.shapes.add_textbox(x + Inches(1.1), y + Inches(0.32), Inches(0.7), Inches(0.65))
    num_tb.text_frame.word_wrap = False
    np = num_tb.text_frame.paragraphs[0]
    np.alignment = PP_ALIGN.CENTER
    nr = np.add_run()
    nr.text = a["num"]
    nr.font.size = Pt(24)
    nr.font.color.rgb = WHITE
    nr.font.bold = True
    nr.font.name = "Calibri"

    # Title
    txt(s4, x + Inches(0.2), y + Inches(1.15), col_w - Inches(0.4), Inches(0.7),
        a["title"], sz=16, clr=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    txt(s4, x + Inches(0.2), y + Inches(1.95), col_w - Inches(0.4), Inches(0.3),
        a["subtitle"], sz=9, clr=MED_GRAY, align=PP_ALIGN.CENTER)

    # Divider
    shape(s4, x + Inches(0.8), y + Inches(2.35), col_w - Inches(1.6), Inches(0.02), RGBColor(0xEE, 0xEE, 0xEE))

    # Bullet points
    bp_y = y + Inches(2.55)
    for j, point in enumerate(a["points"]):
        dot = circle(s4, x + Inches(0.3), bp_y + j * Inches(0.6) + Inches(0.06), Inches(0.1), a["color"])
        txt(s4, x + Inches(0.5), bp_y + j * Inches(0.6), col_w - Inches(0.7), Inches(0.55),
            point, sz=9.5, clr=DARK_GRAY)

    # Priority badge at bottom
    badge_y = y + col_h - Inches(0.55)
    shape(s4, x + Inches(0.4), badge_y, col_w - Inches(0.8), Inches(0.35), a["pri_color"], rounding=True)
    txt(s4, x + Inches(0.4), badge_y + Inches(0.02), col_w - Inches(0.8), Inches(0.33),
        a["priority"], sz=10, clr=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Outlook bar at bottom
shape(s4, Inches(0), Inches(6.3), SW, Inches(1.1), SAMSUNG_DARK)

txt(s4, Inches(0.6), Inches(6.4), Inches(12), Inches(0.35),
    "2026–2028 INDUSTRY OUTLOOK", sz=12, clr=SAMSUNG_ACCENT, bold=True)

outlooks = [
    "Agentic AI becomes the standard for premium cooking",
    "On-device LLMs proliferate for privacy",
    "Matter erodes proprietary ecosystem lock-in",
    "Smart cookware emerges as new product category",
    "Generative recipe AI becomes table stakes",
]

outlook_box = s4.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12), Inches(0.5))
outlook_box.text_frame.word_wrap = True
op = outlook_box.text_frame.paragraphs[0]
for i, o in enumerate(outlooks):
    sep = "    ·    " if i > 0 else ""
    r = op.add_run()
    r.text = f"{sep}{o}"
    r.font.size = Pt(9.5)
    r.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)
    r.font.name = "Calibri"


# ── Save ──
out = "/home/chachapranu/ashu/git/product-research/claude/Samsung_AI_Kitchen_Visual_2026.pptx"
prs.save(out)
print(f"Saved: {out}")
